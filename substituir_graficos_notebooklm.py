# -*- coding: utf-8 -*-
"""
substituir_graficos_notebooklm.py

Abre Oil_Shocks_and_Brazilian_Inflation.pptx e cola as imagens corretas
do TCC em cima dos graficos errados do NotebookLM.

Slide 4: cobre ilustracao de onda com grafico real IPCA Transportes
Slide 5: cobre as 3 barras horizontais com 3 graficos reais (diesel, gas_ref, gas_c)
Slide 7: cobre graficos genericos com os 3 graficos reais corretos
Slide 9: cobre "Hignal Chart" com SD_LP_ipca_transporte_mensal.png
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pathlib import Path

# ─── Caminhos ───────────────────────────────────────────────────────────────
BASE    = Path(r"C:\Users\pedro\OneDrive\Documentos\TCC_python")
LIMPOS  = BASE / "output_petroleo_lp_modelo10_kilian" / "graficos_ols_limpos"
DUPLA   = BASE / "output_petroleo_lp_modelo10_kilian" / "graficos_ols_dupla_faixa"
SD_DIR  = BASE / "output_petroleo_lp_state_dependent_Modelo12"

INPUT  = BASE / "Oil_Shocks_and_Brazilian_Inflation.pptx"
OUTPUT = BASE / "Oil_Shocks_Graficos_Corretos.pptx"

# O slide do NotebookLM tem 17.78" x 9.92" (formato widescreen nativo)
SLIDE_W_IN = 17.78
SLIDE_H_IN = 9.92

def ins(slide, img_path, left_in, top_in, width_in, height_in):
    """Insere imagem com coordenadas em polegadas."""
    p = Path(img_path)
    if not p.exists():
        print(f"  [AVISO] Nao encontrado: {p.name}")
        return
    slide.shapes.add_picture(
        str(p),
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in)
    )
    print(f"  OK: {p.name}  ({left_in}\", {top_in}\", {width_in}\"x{height_in}\")")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — "A Lente Metodologica: Local Projections"
# Gráfico errado: ilustracao decorativa de onda (metade direita do slide)
# Gráfico correto: LP_OLS_Kilian_dupla_faixa_ipca_transporte.png
# Posicao da ilustracao de onda: aprox. x=5.5" y=1.5" w=11.5" h=7.5"
# ════════════════════════════════════════════════════════════════════════════
def slide4(slide):
    print("\nSlide 4 — Substituindo ilustracao de onda por grafico IPCA Transportes...")
    ins(slide,
        DUPLA / "LP_OLS_Kilian_dupla_faixa_ipca_transporte.png",
        left_in=5.4, top_in=1.4,
        width_in=11.8, height_in=7.8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — "Passo 1: A Chegada a Bomba"
# Gráficos errados: 3 barras horizontais (Diesel, Gasolina A, Gasolina C)
# Gráficos corretos: os 3 graficos limpos correspondentes
#
# Layout do slide (17.78" x 9.92"):
#   Titulo + subtitulo ocupa y=0 ate y~2.0"
#   Barra Diesel:      y≈2.1"  h≈2.2"  x≈1.5" w≈15.5"
#   Barra Gasolina A:  y≈4.4"  h≈2.2"  x≈1.5" w≈15.5"
#   Barra Gasolina C:  y≈6.7"  h≈2.5"  x≈1.5" w≈15.5"
#
# Vamos cobrir a area inteira das 3 barras com os 3 graficos lado a lado
# ou empilhados verticalmente replicando o layout original.
# ════════════════════════════════════════════════════════════════════════════
def slide5(slide):
    print("\nSlide 5 — Substituindo barras horizontais por graficos reais...")

    g_w = 4.8   # largura de cada grafico
    g_h = 2.3   # altura de cada grafico
    x_start = 2.0  # inicio horizontal (apos os icones/labels)
    gap = 0.15  # espaco entre graficos na horizontal

    # Linha 1: Diesel (cobre a barra do Diesel)
    # y da barra Diesel aprox 2.0 a 4.3"
    ins(slide,
        LIMPOS / "LP_OLS_Kilian_limpo_diesel.png",
        left_in=x_start, top_in=2.05,
        width_in=15.2, height_in=2.2)

    # Linha 2: Gasolina A (cobre barra Gasolina A)
    # y da barra Gasolina A aprox 4.35 a 6.55"
    ins(slide,
        LIMPOS / "LP_OLS_Kilian_limpo_gasolina_refinaria.png",
        left_in=x_start, top_in=4.3,
        width_in=15.2, height_in=2.2)

    # Linha 3: Gasolina C (cobre barra Gasolina C)
    # y da barra Gasolina C aprox 6.6 a 9.5"
    ins(slide,
        LIMPOS / "LP_OLS_Kilian_limpo_gasolina.png",
        left_in=x_start, top_in=6.55,
        width_in=15.2, height_in=2.6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — "A Diluicao Macro"
# Graficos errados:
#   Esquerda: "A Febre Setorial (Transportes)" — curva generica cinza
#   Direita superior: "Gasolina C vs IPCA Geral" — curva generica
#   Direita inferior: "Diesel vs IPCA Geral"    — curva generica
#
# Layout (17.78" x 9.92"):
#   Titulo/subtitulo: y=0 a y~2.0"
#   Grafico esquerdo (Transportes): x=0.3" y=2.0" w=7.5" h=7.5"
#   Grafico direito superior (Geral): x=8.2" y=2.0" w=9.0" h=3.5"
#   Grafico direito inferior (Diesel geral): x=8.2" y=5.5" w=9.0" h=3.5"
# ════════════════════════════════════════════════════════════════════════════
def slide7(slide):
    print("\nSlide 7 — Substituindo graficos de IPCA Transportes e IPCA Geral...")

    # Grafico esquerdo grande: Gasolina C → IPCA Transportes
    ins(slide,
        LIMPOS / "LP_OLS_Kilian_limpo_ipca_transporte.png",
        left_in=0.3, top_in=2.05,
        width_in=7.8, height_in=7.5)

    # Grafico direito superior: Gasolina C → IPCA Geral
    ins(slide,
        LIMPOS / "LP_OLS_Kilian_limpo_ipca_geral.png",
        left_in=8.2, top_in=2.05,
        width_in=9.1, height_in=3.4)

    # Grafico direito inferior: Diesel → IPCA Geral
    ins(slide,
        LIMPOS / "LP_OLS_Kilian_limpo_diesel.png",
        left_in=8.2, top_in=5.55,
        width_in=9.1, height_in=3.4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — "O Climax: A Valvula da Petrobras"
# Gráfico errado: "Hignal Chart — Reconstruction of Figure 7"
#   Posicao: x~0.5" y~2.0" w~6.5" h~7.5" (metade esquerda do slide)
# Gráfico correto: SD_LP_ipca_transporte_mensal.png
#   Este grafico ja tem as duas curvas (azul pre-2016, vermelho pos-2016)
# ════════════════════════════════════════════════════════════════════════════
def slide9(slide):
    print("\nSlide 9 — Substituindo Hignal Chart pelo SD_LP IPCA Transportes real...")
    ins(slide,
        SD_DIR / "SD_LP_ipca_transporte_mensal.png",
        left_in=0.35, top_in=1.9,
        width_in=7.6, height_in=7.7)


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════
def main():
    print(f"Abrindo: {INPUT}")
    prs = Presentation(str(INPUT))
    print(f"Total de slides: {len(prs.slides)}")
    print(f"Tamanho do slide: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")

    slides = prs.slides

    slide4(slides[3])   # slide 4 = indice 3
    slide5(slides[4])   # slide 5 = indice 4
    slide7(slides[6])   # slide 7 = indice 6
    slide9(slides[8])   # slide 9 = indice 8

    prs.save(str(OUTPUT))
    print(f"\n✅ Salvo em: {OUTPUT}")


if __name__ == "__main__":
    main()
