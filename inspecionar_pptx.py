# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches

prs = Presentation(r'C:\Users\pedro\OneDrive\Documentos\TCC_python\Oil_Shocks_and_Brazilian_Inflation.pptx')
print(f'Total de slides: {len(prs.slides)}')
print()
for i, slide in enumerate(prs.slides):
    has_image = False
    print(f'--- Slide {i+1} ---')
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            l = round(shape.left/914400, 2)
            t = round(shape.top/914400, 2)
            w = round(shape.width/914400, 2)
            h = round(shape.height/914400, 2)
            print(f'  IMG idx={shape.shape_id}: left={l}" top={t}" w={w}" h={h}"')
            has_image = True
        if hasattr(shape, 'text') and shape.text.strip():
            print(f'  TXT: {shape.text.strip()[:100]}')
    if not has_image:
        print('  (sem imagens)')
    print()
