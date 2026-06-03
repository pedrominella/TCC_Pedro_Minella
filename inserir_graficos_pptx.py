# -*- coding: utf-8 -*-
"""
inserir_graficos_pptx.py

Abre o Apresentacao_TCC_Pedro_Minella.pptx e insere os graficos reais
do TCC nos slides corretos, com alinhamento preciso.

Mapeamento de slides (1-indexado):
  Slide 7  -> Petroleo -> Combustiveis: graficos limpos de diesel, gasolina, refinaria, etanol
  Slide 8  -> Gasolina -> IPCA Transp vs Geral: dupla faixa gasolina e grafico limpo IPCA Transp
  Slide 9  -> Diesel impacto temporario: grafico limpo diesel e grafico IPCA transporte diesel
  Slide 11 -> Regimes PPI: SD_LP gasolina, diesel, IPCA transporte por regime
  Slide 12 -> Robustez IV: grafico comparativo Kilian IPCA transporte
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import shutil

# ─────────────────────────────────────────────
# CAMINHOS
# ─────────────────────────────────────────────
BASE_DIR   = Path(r"C:\Users\pedro\OneDrive\Documentos\TCC_python")
KILIAN_DIR = BASE_DIR / "output_petroleo_lp_modelo10_kilian"
LIMPOS_DIR = KILIAN_DIR / "graficos_ols_limpos"
DUPLA_DIR  = KILIAN_DIR / "graficos_ols_dupla_faixa"
SD_DIR     = BASE_DIR / "output_petroleo_lp_state_dependent_Modelo12"

INPUT_PPTX  = BASE_DIR / "Apresentacao_TCC_Pedro_Minella.pptx"
OUTPUT_PPTX = BASE_DIR / "Apresentacao_TCC_Pedro_Minella.pptx"

# Cores de fundo para labels das figuras
NAVY      = RGBColor(0x10, 0x1C, 0x3E)
DARK_BLUE = RGBColor(0x1A, 0x2D, 0x5A)
ACCENT    = RGBColor(0x2C, 0x82, 0xC9)
GOLD      = RGBColor(0xF5, 0xC5, 0x18)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xD6, 0xE0, 0xF0)
MID_GRAY  = RGBColor(0x8A, 0xA3, 0xC4)
GREEN_OK  = RGBColor(0x2E, 0xCC, 0x71)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_picture_safe(slide, img_path, left, top, width, height):
    """Insere imagem no slide com tratamento de erro."""
    path = Path(img_path)
    if not path.exists():
        print(f"  [AVISO] Imagem nao encontrada: {path}")
        return None
    pic = slide.shapes.add_picture(str(path), left, top, width, height)
    return pic


def add_label_box(slide, text, left, top, width, height,
                  bg_color=DARK_BLUE, txt_color=ACCENT, font_size=Pt(10),
                  bold=False, align=PP_ALIGN.CENTER):
    """Adiciona uma caixa de texto/label sobre ou abaixo de uma imagem."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = txt_color
    return shape


def add_rect_bg(slide, left, top, width, height, fill_color=DARK_BLUE):
    """Adiciona retangulo de fundo para grupo de graficos."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ═══════════════════════════════════════════════════════
# SLIDE 7 — Petroleo → Combustiveis
# Substitui o painel esquerdo (tabela manual) por 4 graficos:
#   Gasolina Refinaria | Gasolina Bomba
#   Diesel             | Etanol
# ═══════════════════════════════════════════════════════
def inserir_slide7(slide):
    print("  Inserindo graficos no Slide 7 (Petroleo -> Combustiveis)...")

    # Limpa o painel esquerdo da tabela manual (shapes com texto de dados)
    # Os graficos ficam na metade esquerda (0 ate ~7.8 polegadas)
    # e o painel de insights (azul) permanece na direita (>= 8.0)
    shapes_to_delete = []
    for shape in slide.shapes:
        l = shape.left or 0
        # Removemos shapes textuais que estao na area da tabela (esquerda)
        if hasattr(shape, "text") and l < Inches(7.9) and l > Inches(0.1):
            t = (shape.text or "").strip()
            # Detecta linhas da tabela (h=0, h=1...) e headers
            if (t.startswith("h =") or t.startswith("Horizonte") or
                    t.startswith("Diesel") or t.startswith("Gasolina") or
                    t.startswith("Etanol") or t.startswith("0,") or
                    t.startswith("−") or "***" in t or "**" in t or
                    t.startswith("* p<") or "Coeficientes acumulados" in t):
                shapes_to_delete.append(shape)

    for s in shapes_to_delete:
        sp = s._element
        sp.getparent().remove(sp)

    # Remove retangulos coloridos da tabela (background boxes na area esquerda)
    rects_to_delete = []
    for shape in slide.shapes:
        l = shape.left or 0
        t_ = shape.top or 0
        w_ = shape.width or 0
        h_ = shape.height or 0
        # Shapes pequenos de celula de tabela na area esquerda
        if (l > Inches(0.1) and l < Inches(7.9) and
                t_ > Inches(1.5) and t_ < Inches(6.3) and
                w_ < Inches(2.5) and h_ < Inches(0.6)):
            if not hasattr(shape, "text") or not (shape.text or "").strip():
                rects_to_delete.append(shape)

    for s in rects_to_delete:
        sp = s._element
        sp.getparent().remove(sp)

    # Area disponivel: esquerda (0.3 a 7.9), topo 1.6 a 6.2
    # Dividimos em grade 2x2 de graficos
    # Cada grafico: ~3.5 x 2.1 polegadas
    g_w = Inches(3.65)
    g_h = Inches(2.15)
    x1 = Inches(0.3)
    x2 = Inches(4.1)
    y1 = Inches(1.62)
    y2 = Inches(3.88)

    # Fundo unificado para a grade de graficos
    add_rect_bg(slide, left=Inches(0.3), top=Inches(1.55),
                width=Inches(7.5), height=Inches(4.85),
                fill_color=DARK_BLUE)

    # Grafico 1: Gasolina Refinaria (topo esquerda)
    add_picture_safe(slide,
        LIMPOS_DIR / "LP_OLS_Kilian_limpo_gasolina_refinaria.png",
        left=x1, top=y1, width=g_w, height=g_h)
    add_label_box(slide, "Gasolina A (Refinaria)",
                  left=x1, top=y1 - Inches(0.28), width=g_w, height=Inches(0.26),
                  txt_color=ACCENT, font_size=Pt(9.5), bold=True)

    # Grafico 2: Gasolina Bomba (topo direita)
    add_picture_safe(slide,
        LIMPOS_DIR / "LP_OLS_Kilian_limpo_gasolina.png",
        left=x2, top=y1, width=g_w, height=g_h)
    add_label_box(slide, "Gasolina C (Bomba ao Consumidor)",
                  left=x2, top=y1 - Inches(0.28), width=g_w, height=Inches(0.26),
                  txt_color=GREEN_OK, font_size=Pt(9.5), bold=True)

    # Grafico 3: Diesel (base esquerda)
    add_picture_safe(slide,
        LIMPOS_DIR / "LP_OLS_Kilian_limpo_diesel.png",
        left=x1, top=y2, width=g_w, height=g_h)
    add_label_box(slide, "Oleo Diesel",
                  left=x1, top=y2 - Inches(0.28), width=g_w, height=Inches(0.26),
                  txt_color=GOLD, font_size=Pt(9.5), bold=True)

    # Grafico 4: Etanol (base direita)
    add_picture_safe(slide,
        LIMPOS_DIR / "LP_OLS_Kilian_limpo_etanol.png",
        left=x2, top=y2, width=g_w, height=g_h)
    add_label_box(slide, "Etanol",
                  left=x2, top=y2 - Inches(0.28), width=g_w, height=Inches(0.26),
                  txt_color=MID_GRAY, font_size=Pt(9.5), bold=True)

    # Nota de fonte
    add_label_box(slide,
                  "Fonte: Elaboracao propria. Modelo LP-OLS com Controle Kilian. IC 90% e 95% (areas sombreadas). Eixo Y = p.p. acumulados.",
                  left=Inches(0.3), top=Inches(6.12),
                  width=Inches(7.5), height=Inches(0.25),
                  txt_color=MID_GRAY, font_size=Pt(8), bold=False)


# ═══════════════════════════════════════════════════════
# SLIDE 8 — Gasolina → IPCA Transportes vs Geral
# Substitui as colunas de dados por dois graficos:
#   Esquerda: LP_OLS_Kilian_dupla_faixa_ipca_transporte (destaque: 43%)
#   Direita:  LP_OLS_Kilian_limpo_ipca_geral
# ═══════════════════════════════════════════════════════
def inserir_slide8(slide):
    print("  Inserindo graficos no Slide 8 (Gasolina -> IPCA Transp vs Geral)...")

    # Remove shapes textuais das duas colunas de dados
    shapes_to_delete = []
    for shape in slide.shapes:
        l = shape.left or 0
        t_ = shape.top or 0
        w_ = shape.width or 0
        if not hasattr(shape, "text"):
            continue
        txt = (shape.text or "").strip()
        # Remove linhas de dados (coeficientes, h=X, textos de dados)
        if t_ > Inches(1.5) and t_ < Inches(6.1):
            if (txt.startswith("h =") or txt.startswith("0,") or
                    txt.startswith("−") or "***" in txt or "**" in txt or
                    "repasse" in txt.lower() or "aceleração" in txt.lower() or
                    "consolidação" in txt.lower() or "estabilização" in txt.lower() or
                    "mantém" in txt.lower() or "persistente" in txt.lower() or
                    "efeito inicial" in txt.lower() or "pico máximo" in txt.lower() or
                    "sem significância" in txt.lower() or
                    "Pass-through" in txt):
                shapes_to_delete.append(shape)

    # Remove também os retangulos de destaque de dados (barras coloridas)
    for shape in slide.shapes:
        l = shape.left or 0
        t_ = shape.top or 0
        w_ = shape.width or 0
        h_ = shape.height or 0
        txt = ""
        if hasattr(shape, "text"):
            txt = (shape.text or "").strip()
        if (t_ > Inches(5.2) and t_ < Inches(5.7) and
                (w_ > Inches(5.0)) and "Pass-through" not in txt):
            shapes_to_delete.append(shape)

    for s in shapes_to_delete:
        try:
            sp = s._element
            sp.getparent().remove(sp)
        except Exception:
            pass

    # Remove headers IPCA TRANSPORTES e IPCA GERAL preservando o layout geral
    headers_to_delete = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        txt = (shape.text or "").strip()
        if "IPCA TRANSPORTES" in txt or "IPCA GERAL" in txt:
            headers_to_delete.append(shape)
    for s in headers_to_delete:
        try:
            sp = s._element
            sp.getparent().remove(sp)
        except Exception:
            pass

    # Grafico ESQUERDA: IPCA Transportes com dupla faixa (IC 90% e 95%)
    # Area: x=0.3 a 6.4, y=1.55 a 6.1
    g_w_l = Inches(5.9)
    g_h_l = Inches(4.0)
    x_left = Inches(0.3)
    y_graf = Inches(1.62)

    add_rect_bg(slide, left=Inches(0.3), top=Inches(1.55),
                width=Inches(6.1), height=Inches(4.7), fill_color=DARK_BLUE)

    add_label_box(slide, "IPCA TRANSPORTES  ✅  (LP-OLS Kilian — IC 90% e 95%)",
                  left=Inches(0.35), top=Inches(1.58),
                  width=Inches(6.0), height=Inches(0.3),
                  txt_color=GREEN_OK, font_size=Pt(10), bold=True)

    add_picture_safe(slide,
        DUPLA_DIR / "LP_OLS_Kilian_dupla_faixa_ipca_transporte.png",
        left=x_left, top=y_graf + Inches(0.3),
        width=g_w_l, height=g_h_l)

    # Caixa destaque 43%
    add_rect_bg(slide, left=Inches(0.3), top=Inches(5.7),
                width=Inches(6.1), height=Inches(0.55),
                fill_color=RGBColor(0x0D, 0x3D, 0x20))
    add_label_box(slide,
                  "Pass-through de longo prazo: ~43%  |  t-stat > 3,6 em todo o horizonte",
                  left=Inches(0.35), top=Inches(5.73),
                  width=Inches(6.0), height=Inches(0.45),
                  txt_color=GREEN_OK, font_size=Pt(10.5), bold=True, align=PP_ALIGN.CENTER)

    # Grafico DIREITA: IPCA Geral (limpo)
    add_rect_bg(slide, left=Inches(6.9), top=Inches(1.55),
                width=Inches(6.1), height=Inches(4.7), fill_color=DARK_BLUE)

    add_label_box(slide, "IPCA GERAL  ⚠️  (LP-OLS Kilian — IC 90% e 95%)",
                  left=Inches(6.95), top=Inches(1.58),
                  width=Inches(6.0), height=Inches(0.3),
                  txt_color=MID_GRAY, font_size=Pt(10), bold=True)

    add_picture_safe(slide,
        DUPLA_DIR / "LP_OLS_Kilian_dupla_faixa_ipca_geral.png",
        left=Inches(6.9), top=y_graf + Inches(0.3),
        width=g_w_l, height=g_h_l)

    add_rect_bg(slide, left=Inches(6.9), top=Inches(5.7),
                width=Inches(6.1), height=Inches(0.55),
                fill_color=RGBColor(0x2A, 0x1A, 0x0A))
    add_label_box(slide,
                  "Pass-through diluido: < 0,06 p.p.  |  Insignificante apos o 6.o mes",
                  left=Inches(6.95), top=Inches(5.73),
                  width=Inches(6.0), height=Inches(0.45),
                  txt_color=MID_GRAY, font_size=Pt(10.5), bold=True, align=PP_ALIGN.CENTER)

    # Nota de fonte abaixo
    add_label_box(slide,
                  "Fonte: Elaboracao propria. Gasolina C (Bomba) como choque. Controles: cambio, IBC-Br, Selic, Focus, Kilian. HAC Newey-West.",
                  left=Inches(0.3), top=Inches(6.33),
                  width=Inches(12.7), height=Inches(0.22),
                  txt_color=MID_GRAY, font_size=Pt(8), bold=False)


# ═══════════════════════════════════════════════════════
# SLIDE 9 — Diesel: Choque Temporario
# Substitui o grafico de barras manual por 2 graficos reais:
#   Grafico principal: LP_OLS_Kilian_dupla_faixa_diesel (IPCA Transportes)
#   Grafico auxiliar:  LP_OLS_Kilian_limpo_ipca_transporte (petroleo->transp.)
# ═══════════════════════════════════════════════════════
def inserir_slide9(slide):
    print("  Inserindo graficos no Slide 9 (Diesel - Impacto Temporario)...")

    # Remove o grafico de barras manual (retangulos de barra)
    shapes_to_delete = []
    for shape in slide.shapes:
        l = shape.left or 0
        t_ = shape.top or 0
        w_ = shape.width or 0
        h_ = shape.height or 0

        # Retangulos pequenos das barras manuais (area de grafico esquerda)
        if (l > Inches(0.35) and l < Inches(8.0) and
                t_ > Inches(1.6) and t_ < Inches(6.5) and
                w_ < Inches(0.7) and h_ < Inches(2.5)):
            shapes_to_delete.append(shape)

        # Labels das barras (h=0, h=1... e valores)
        if hasattr(shape, "text"):
            txt = (shape.text or "").strip()
            if (txt.startswith("h=") or (len(txt) > 0 and txt[0] in "0123456789−-") or
                    "Diesel" in txt and "IPCA" in txt):
                if l < Inches(8.0) and t_ > Inches(1.5) and t_ < Inches(6.7):
                    shapes_to_delete.append(shape)

    # Linha zero
    for shape in slide.shapes:
        l = shape.left or 0
        t_ = shape.top or 0
        h_ = shape.height or 0
        w_ = shape.width or 0
        if (l > Inches(0.3) and l < Inches(1.0) and
                t_ > Inches(3.0) and t_ < Inches(4.5) and
                h_ < Pt(3) and w_ > Inches(5.0)):
            shapes_to_delete.append(shape)

    for s in shapes_to_delete:
        try:
            sp = s._element
            sp.getparent().remove(sp)
        except Exception:
            pass

    # Remove blocos de legenda do grafico manual
    leg_to_del = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        txt = (shape.text or "").strip()
        if ("Diesel → IPCA Transportes" in txt or "Diesel → IPCA Geral" in txt):
            l = shape.left or 0
            if l < Inches(8.0):
                leg_to_del.append(shape)
    for s in leg_to_del:
        try:
            sp = s._element
            sp.getparent().remove(sp)
        except Exception:
            pass

    # Grafico principal DIESEL -> IPCA TRANSPORTES (dupla faixa, ocupa 2/3 da esquerda)
    g_w_main = Inches(7.4)
    g_h_main = Inches(4.5)

    add_rect_bg(slide, left=Inches(0.3), top=Inches(1.55),
                width=Inches(7.6), height=Inches(5.0), fill_color=DARK_BLUE)

    add_label_box(slide, "Diesel → IPCA Transportes  |  IC 90% e 95%",
                  left=Inches(0.35), top=Inches(1.6),
                  width=Inches(7.5), height=Inches(0.28),
                  txt_color=GREEN_OK, font_size=Pt(10), bold=True)

    add_picture_safe(slide,
        DUPLA_DIR / "LP_OLS_Kilian_dupla_faixa_diesel.png",
        left=Inches(0.3), top=Inches(1.92),
        width=g_w_main, height=g_h_main)

    add_label_box(slide,
                  "Fonte: Elaboracao propria. LP-OLS com Controle Kilian. Choque = 1 p.p. no Oleo Diesel. HAC Newey-West.",
                  left=Inches(0.3), top=Inches(6.5),
                  width=Inches(7.6), height=Inches(0.22),
                  txt_color=MID_GRAY, font_size=Pt(8))


# ═══════════════════════════════════════════════════════
# SLIDE 11 — Regimes PPI: Antes e Depois de 2016
# Insere os 3 graficos State-Dependent LP:
#   Topo esquerda:  SD_LP_dln_gasolina.png
#   Topo direita:   SD_LP_dln_diesel.png
#   Base centro:    SD_LP_ipca_transporte_mensal.png (o mais importante)
# ═══════════════════════════════════════════════════════
def inserir_slide11(slide):
    print("  Inserindo graficos no Slide 11 (Regimes PPI)...")

    # Remove a tabela de dados (blocos de texto numericos e headers)
    shapes_to_delete = []
    for shape in slide.shapes:
        l = shape.left or 0
        t_ = shape.top or 0
        if not hasattr(shape, "text"):
            continue
        txt = (shape.text or "").strip()
        # Celulas da tabela de regime
        if (t_ > Inches(1.5) and t_ < Inches(5.7) and l < Inches(7.7)):
            if (txt.startswith("h =") or txt.startswith("0,") or
                    txt.startswith("[") or "n.sig." in txt or
                    txt.startswith("Horizonte") or txt.startswith("Pré") or
                    txt.startswith("Pós") or txt.startswith("IC") or
                    txt.startswith("Wald") or txt.startswith("< 0") or
                    "Linha AZUL" in txt or "Linha VERMELHA" in txt or
                    "coeficiente" in txt.lower()):
                shapes_to_delete.append(shape)

    # Remove retangulos de header da tabela
    for shape in slide.shapes:
        l = shape.left or 0
        t_ = shape.top or 0
        w_ = shape.width or 0
        h_ = shape.height or 0
        txt = ""
        if hasattr(shape, "text"):
            txt = (shape.text or "").strip()
        if (l < Inches(7.5) and t_ > Inches(1.5) and t_ < Inches(6.0) and
                w_ < Inches(1.7) and h_ < Inches(0.55) and not txt):
            shapes_to_delete.append(shape)

    for s in shapes_to_delete:
        try:
            sp = s._element
            sp.getparent().remove(sp)
        except Exception:
            pass

    # Remocao dos 2 boxes coloridos de destaque Azul/Vermelho (linhas de regime)
    destaque_del = []
    for shape in slide.shapes:
        l = shape.left or 0
        t_ = shape.top or 0
        w_ = shape.width or 0
        h_ = shape.height or 0
        if (t_ > Inches(5.2) and t_ < Inches(6.5) and
                l > Inches(0.1) and l < Inches(7.5) and
                w_ > Inches(5.0)):
            destaque_del.append(shape)
    for s in destaque_del:
        try:
            sp = s._element
            sp.getparent().remove(sp)
        except Exception:
            pass

    # Agora insere os 3 graficos de regimes
    # Layout:
    #   [ Gasolina por Regime ]  |  [ painel de insights (direita, mantem) ]
    #   [ Diesel por Regime   ]  |
    #   [  IPCA Transp Regime (grande, central) ]

    # Fundo para area de graficos
    add_rect_bg(slide, left=Inches(0.3), top=Inches(1.55),
                width=Inches(6.8), height=Inches(4.7), fill_color=DARK_BLUE)

    # Grafico 1: Gasolina por regime (metade superior)
    add_label_box(slide, "Gasolina → Preco (Pre vs. Pos-2016)",
                  left=Inches(0.35), top=Inches(1.58),
                  width=Inches(3.2), height=Inches(0.28),
                  txt_color=ACCENT, font_size=Pt(9), bold=True)

    add_picture_safe(slide,
        SD_DIR / "SD_LP_dln_gasolina.png",
        left=Inches(0.3), top=Inches(1.88),
        width=Inches(3.3), height=Inches(2.1))

    # Grafico 2: Diesel por regime (ao lado do grafico de gasolina)
    add_label_box(slide, "Diesel → Preco (Pre vs. Pos-2016)",
                  left=Inches(3.75), top=Inches(1.58),
                  width=Inches(3.2), height=Inches(0.28),
                  txt_color=GOLD, font_size=Pt(9), bold=True)

    add_picture_safe(slide,
        SD_DIR / "SD_LP_dln_diesel.png",
        left=Inches(3.75), top=Inches(1.88),
        width=Inches(3.3), height=Inches(2.1))

    # Grafico 3: IPCA Transportes por regime (destaque principal - maior)
    add_label_box(slide, "IPCA Transportes: Pre-2016 (azul) vs. Pos-2016 / PPI (vermelho)  — O achado central do TCC",
                  left=Inches(0.35), top=Inches(4.08),
                  width=Inches(6.7), height=Inches(0.28),
                  txt_color=GREEN_OK, font_size=Pt(9.5), bold=True)

    add_picture_safe(slide,
        SD_DIR / "SD_LP_ipca_transporte_mensal.png",
        left=Inches(0.3), top=Inches(4.38),
        width=Inches(6.8), height=Inches(1.85))

    add_label_box(slide,
                  "Fonte: Elaboracao propria. Modelo 12 (LP State-Dependent). Corte: setembro/2016 (PPI Petrobras). IC 90%.",
                  left=Inches(0.3), top=Inches(6.28),
                  width=Inches(6.8), height=Inches(0.22),
                  txt_color=MID_GRAY, font_size=Pt(8))


# ═══════════════════════════════════════════════════════
# SLIDE 12 — Robustez: IV com OSNS
# Adiciona o grafico comparativo Kilian IPCA Transporte
# ao lado esquerdo (acima da cascata de 3 passos)
# ═══════════════════════════════════════════════════════
def inserir_slide12(slide):
    print("  Inserindo graficos no Slide 12 (Robustez IV)...")

    # Nao remove nada do slide 12 (a cascata de 3 boxes e essencial)
    # Apenas insere um grafico comparativo LP-OLS vs LP-IV em area livre

    # Verifica se ha espaco acima da cascata de steps (que comeca em y~1.6)
    # Inserimos o grafico do comparativo Kilian IPCA Transporte
    # em area overlay no topo direito

    add_rect_bg(slide, left=Inches(8.6), top=Inches(1.58),
                width=Inches(4.4), height=Inches(2.6), fill_color=DARK_BLUE)

    add_label_box(slide, "Resposta Acumulada: Petroleo → IPCA Transportes",
                  left=Inches(8.65), top=Inches(1.6),
                  width=Inches(4.3), height=Inches(0.28),
                  txt_color=ACCENT, font_size=Pt(9), bold=True)

    add_picture_safe(slide,
        KILIAN_DIR / "comparativo_Kilian_ipca_transporte.png",
        left=Inches(8.6), top=Inches(1.9),
        width=Inches(4.4), height=Inches(2.2))

    add_label_box(slide,
                  "LP-OLS (azul) vs LP-IV Kilian (laranja). Curvas sobrepostas confirmam: vies de endogeneidade e desprezivel.",
                  left=Inches(8.6), top=Inches(4.12),
                  width=Inches(4.4), height=Inches(0.38),
                  txt_color=MID_GRAY, font_size=Pt(7.5))


# ═══════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════
def main():
    print(f"Abrindo: {INPUT_PPTX}")
    prs = Presentation(str(INPUT_PPTX))

    slides = prs.slides
    print(f"Total de slides: {len(slides)}")

    # Slide 7 (indice 6)
    inserir_slide7(slides[6])

    # Slide 8 (indice 7)
    inserir_slide8(slides[7])

    # Slide 9 (indice 8)
    inserir_slide9(slides[8])

    # Slide 11 (indice 10)
    inserir_slide11(slides[10])

    # Slide 12 (indice 11)
    inserir_slide12(slides[11])

    print(f"\nSalvando em: {OUTPUT_PPTX}")
    prs.save(str(OUTPUT_PPTX))
    print("Apresentacao com graficos salva com sucesso!")


if __name__ == "__main__":
    main()
