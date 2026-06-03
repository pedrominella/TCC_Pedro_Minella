# -*- coding: utf-8 -*-
"""
rebuild_com_graficos.py

Estrategia: usa cada slide do NotebookLM como fundo (ja esta la como imagem),
e adiciona os graficos reais do TCC por cima, com fundo branco e posicao
cuidadosamente calibrada para cobrir exatamente a area do grafico errado.

Slide 4 - "A Lente Metodologica": cobre a ilustracao de onda (lado direito)
Slide 5 - "Passo 1: A Chegada a Bomba": cobre as 3 barras horizontais
Slide 7 - "A Diluicao Macro": cobre os 3 graficos genericos
Slide 9 - "O Climax: A Valvula da Petrobras": cobre o Hignal Chart (lado esq.)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pathlib import Path
from lxml import etree

BASE   = Path(r"C:\Users\pedro\OneDrive\Documentos\TCC_python")
LIMPOS = BASE / "output_petroleo_lp_modelo10_kilian" / "graficos_ols_limpos"
DUPLA  = BASE / "output_petroleo_lp_modelo10_kilian" / "graficos_ols_dupla_faixa"
SD_DIR = BASE / "output_petroleo_lp_state_dependent_Modelo12"

INPUT  = BASE / "Oil_Shocks_and_Brazilian_Inflation.pptx"
OUTPUT = BASE / "Oil_Shocks_Graficos_Corretos.pptx"

# Slide do NotebookLM: 17.78" x 10.00"
W = 17.78
H = 10.00

def pic(slide, path, l, t, w, h):
    """Insere imagem com fundo branco limpo."""
    p = Path(path)
    if not p.exists():
        print(f"  [FALTANDO] {p.name}")
        return
    # Fundo branco antes da imagem (cobre o grafico errado com retangulo branco)
    bg = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Imagem sobre o fundo branco
    slide.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))
    print(f"  OK: {p.name}")

def label(slide, text, l, t, w, h, size=Pt(11), bold=False,
          color=RGBColor(0x12, 0x34, 0x56), align=PP_ALIGN.LEFT, bg=None):
    """Adiciona texto com fundo opcional."""
    if bg:
        r = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        r.fill.solid()
        r.fill.fore_color.rgb = bg
        r.line.fill.background()
        r.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.alignment = align
    run = p0.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — "A Lente Metodologica: Local Projections"
#
# Layout original:
#   - Esquerda (0 a ~8"): texto com bullets (OK, manter)
#   - Direita (~8" a 17.78"): ilustracao de onda decorativa (SUBSTITUIR)
#
# Substituicao: grafico IPCA Transportes com dupla faixa
# O grafico cobre a area da onda. Adiciona legenda explicativa acima.
# ══════════════════════════════════════════════════════════════════════
def fix_slide4(slide):
    print("\nSlide 4...")
    # Cobre a area inteira da ilustracao de onda (metade direita)
    # Deixa margem para o titulo e bullets nao serem afetados
    pic(slide, DUPLA / "LP_OLS_Kilian_dupla_faixa_ipca_transporte.png",
        l=7.9, t=1.5, w=9.5, h=7.8)

    # Label de titulo do grafico (sobre fundo azul escuro, estilo do template)
    label(slide,
          "Funcao de Resposta Acumulada: Gasolina C -> IPCA Transportes",
          l=7.95, t=1.52, w=9.4, h=0.45,
          size=Pt(13), bold=True,
          color=RGBColor(0xFF, 0xFF, 0xFF),
          align=PP_ALIGN.LEFT,
          bg=RGBColor(0x12, 0x3A, 0x5E))

    # Nota de rodape do grafico
    label(slide,
          "Modelo LP-OLS com Controle Kilian. IC 90% (cinza escuro) e 95% (cinza claro). Eixo Y = p.p. acumulados por 1 p.p. no preco da Gasolina C.",
          l=7.95, t=9.08, w=9.4, h=0.45,
          size=Pt(8.5), bold=False,
          color=RGBColor(0x55, 0x55, 0x55),
          align=PP_ALIGN.LEFT,
          bg=RGBColor(0xF0, 0xF0, 0xF0))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — "Passo 1: A Chegada a Bomba"
#
# Layout original (17.78" x 10"):
#   Titulo: y=0 a ~1.6"
#   Barra Diesel:       y~1.6"  a  y~4.2"   (h~2.6")
#   Barra Gasolina A:   y~4.2"  a  y~6.7"   (h~2.5")
#   Barra Gasolina C:   y~6.7"  a  y~9.5"   (h~2.8")
#   Os icones ficam na esquerda (x=0 a ~2.5")
#   As barras comecam em x~2.5" e vao ate x~16.5"
#
# Substituicao: 3 graficos LP limpos empilhados cobrindo as barras
# ══════════════════════════════════════════════════════════════════════
def fix_slide5(slide):
    print("\nSlide 5...")

    # DIESEL (linha 1)
    pic(slide, LIMPOS / "LP_OLS_Kilian_limpo_diesel.png",
        l=2.4, t=1.65, w=14.8, h=2.5)
    label(slide, "Diesel  |  Resposta acumulada ao choque do Brent (h=0 a 12 meses)",
          l=2.4, t=1.67, w=14.8, h=0.38,
          size=Pt(11), bold=True,
          color=RGBColor(0xFF, 0xFF, 0xFF),
          bg=RGBColor(0x18, 0x3A, 0x5C))

    # GASOLINA A (linha 2)
    pic(slide, LIMPOS / "LP_OLS_Kilian_limpo_gasolina_refinaria.png",
        l=2.4, t=4.2, w=14.8, h=2.5)
    label(slide, "Gasolina A (Refinaria)  |  Resposta acumulada ao choque do Brent",
          l=2.4, t=4.22, w=14.8, h=0.38,
          size=Pt(11), bold=True,
          color=RGBColor(0xFF, 0xFF, 0xFF),
          bg=RGBColor(0x18, 0x3A, 0x5C))

    # GASOLINA C (linha 3)
    pic(slide, LIMPOS / "LP_OLS_Kilian_limpo_gasolina.png",
        l=2.4, t=6.75, w=14.8, h=2.85)
    label(slide, "Gasolina C (Bomba)  |  Resposta acumulada ao choque do Brent",
          l=2.4, t=6.77, w=14.8, h=0.38,
          size=Pt(11), bold=True,
          color=RGBColor(0xFF, 0xFF, 0xFF),
          bg=RGBColor(0x18, 0x3A, 0x5C))

    # Nota de rodape unica
    label(slide,
          "Fonte: Elaboracao propria. Modelo LP-OLS Kilian (Modelo 10). HAC Newey-West. *** p<0,01  ** p<0,05  * p<0,10",
          l=2.4, t=9.65, w=14.8, h=0.32,
          size=Pt(8.5), bold=False,
          color=RGBColor(0x55, 0x55, 0x55),
          bg=RGBColor(0xF4, 0xF4, 0xF4))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — "A Diluicao Macro"
#
# Layout original:
#   Titulo+subtitulo: y=0 a ~1.8"
#   Grafico esquerdo grande (Transportes): x=0.3" y=1.9" w=7.8" h=7.8"
#   Grafico direito superior (Geral):      x=8.2" y=1.9" w=9.0" h=3.5"
#   Grafico direito inferior (Diesel):     x=8.2" y=5.5" w=9.0" h=3.5"
#   Box de conclusao: y~8.9" (manter, nao e grafico)
#
# Substituicao: 3 graficos reais nas mesmas posicoes
# ══════════════════════════════════════════════════════════════════════
def fix_slide7(slide):
    print("\nSlide 7...")

    # Grafico esquerdo: IPCA Transportes (grande)
    pic(slide, LIMPOS / "LP_OLS_Kilian_limpo_ipca_transporte.png",
        l=0.25, t=1.85, w=7.9, h=7.85)
    label(slide, "Gasolina C -> IPCA Transportes",
          l=0.25, t=1.87, w=7.9, h=0.42,
          size=Pt(12), bold=True,
          color=RGBColor(0xFF, 0xFF, 0xFF),
          bg=RGBColor(0x7B, 0x1C, 0x1C))   # vermelho escuro do template

    # Grafico direito superior: IPCA Geral
    pic(slide, LIMPOS / "LP_OLS_Kilian_limpo_ipca_geral.png",
        l=8.2, t=1.85, w=9.3, h=3.6)
    label(slide, "Gasolina C -> IPCA Geral",
          l=8.2, t=1.87, w=9.3, h=0.38,
          size=Pt(11), bold=True,
          color=RGBColor(0xFF, 0xFF, 0xFF),
          bg=RGBColor(0x18, 0x3A, 0x5C))

    # Grafico direito inferior: Diesel -> IPCA Geral
    pic(slide, LIMPOS / "LP_OLS_Kilian_limpo_diesel.png",
        l=8.2, t=5.5, w=9.3, h=3.6)
    label(slide, "Diesel -> IPCA Geral",
          l=8.2, t=5.52, w=9.3, h=0.38,
          size=Pt(11), bold=True,
          color=RGBColor(0xFF, 0xFF, 0xFF),
          bg=RGBColor(0x18, 0x3A, 0x5C))

    # Nota de rodape
    label(slide,
          "Fonte: Elaboracao propria. LP-OLS com Controle Kilian. IC 90% e 95%. *** p<0,01",
          l=0.25, t=9.72, w=17.2, h=0.28,
          size=Pt(8.5),
          color=RGBColor(0x55, 0x55, 0x55),
          bg=RGBColor(0xF4, 0xF4, 0xF4))


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — "O Climax: A Valvula da Petrobras"
#
# Layout original:
#   Titulo+subtitulo: y=0 a ~1.7"
#   "Hignal Chart" (grafico inventado): x=0.3" y=1.9" w=7.5" h=7.8"
#   Painel direito (Dois Brasis, Wald, etc.): x=8.2" a 17.78" (MANTER)
#
# Substituicao: SD_LP_ipca_transporte_mensal.png no lugar do Hignal Chart
# ══════════════════════════════════════════════════════════════════════
def fix_slide9(slide):
    print("\nSlide 9...")

    # Cobre o "Hignal Chart" com o grafico real de regimes
    pic(slide, SD_DIR / "SD_LP_ipca_transporte_mensal.png",
        l=0.2, t=1.85, w=7.8, h=7.85)

    # Label do grafico real (estilo do template - caixa azul escura)
    label(slide,
          "IPCA Transportes: Pre-2016 (azul) vs. Pos-2016/PPI (vermelho)",
          l=0.2, t=1.87, w=7.8, h=0.42,
          size=Pt(11), bold=True,
          color=RGBColor(0xFF, 0xFF, 0xFF),
          bg=RGBColor(0x12, 0x3A, 0x5E))

    # Nota de fonte
    label(slide,
          "Fonte: Elaboracao propria. Modelo LP State-Dependent (Modelo 12). Corte: set/2016 (PPI). IC 90%. Teste de Wald p<0,05.",
          l=0.2, t=9.72, w=7.8, h=0.28,
          size=Pt(8.5),
          color=RGBColor(0x55, 0x55, 0x55),
          bg=RGBColor(0xF4, 0xF4, 0xF4))


# ══════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════
def main():
    print(f"Abrindo: {INPUT}")
    prs = Presentation(str(INPUT))
    print(f"Slides: {len(prs.slides)}  |  Tamanho: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")

    fix_slide4(prs.slides[3])
    fix_slide5(prs.slides[4])
    fix_slide7(prs.slides[6])
    fix_slide9(prs.slides[8])

    prs.save(str(OUTPUT))
    print(f"\nSalvo em: {OUTPUT}")


if __name__ == "__main__":
    main()
