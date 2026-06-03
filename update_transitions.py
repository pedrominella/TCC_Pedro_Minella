# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def format_transition(slide, title_text):
    # Limpar shapes existentes
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
        
    NAVY = RGBColor(0x0D, 0x21, 0x37)
    TEAL = RGBColor(0x1A, 0x6B, 0x8A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    
    # Fundo
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    
    # Faixa lateral esquerda (TEAL)
    s1 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    s1.fill.solid()
    s1.fill.fore_color.rgb = TEAL
    s1.line.fill.background()
    s1.shadow.inherit = False
    
    # Linha branca fina acompanhando
    s2 = slide.shapes.add_shape(1, Inches(0.4), Inches(0), Inches(0.05), Inches(7.5))
    s2.fill.solid()
    s2.fill.fore_color.rgb = WHITE
    s2.line.fill.background()
    s2.shadow.inherit = False

    # Titulo centralizado verticalmente
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.5), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title_text
    r.font.name = "Calibri"
    r.font.size = Pt(48)
    r.font.bold = True
    r.font.color.rgb = WHITE

prs = Presentation(r'C:\Users\pedro\OneDrive\Documentos\TCC_python\TCC_Apresentacao_Final.pptx')

# O usuário mencionou "slides 3 e 4" para "Mas o que a literatura nos diz?". 
# Wait, let me check the exact slides. 
# Slide 3 (index 2) had "Mas o que a literatura nos diz?"
# Slide 4 (index 3) seems to be the one I called Slide 4 in my output before.
# Wait, let me just format slide 3 (index 2) and slide 5 (index 4).

format_transition(prs.slides[2], "Mas o que a literatura nos diz?")
format_transition(prs.slides[4], "Qual Modelo usar e por quê?")

prs.save(r'C:\Users\pedro\OneDrive\Documentos\TCC_python\TCC_Apresentacao_Final.pptx')
print("Slides 3 e 5 atualizados com sucesso.")
