# -*- coding: utf-8 -*-
"""
tcc_novo_template.py
Cria do zero uma apresentacao profissional de 12 slides para defesa do TCC.
Design: clean academico com fundo claro, tipografia moderna, graficos reais.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# ─── Caminhos ───────────────────────────────────────────────────────────────
BASE   = Path(r"C:\Users\pedro\OneDrive\Documentos\TCC_python")
LIMPOS = BASE / "output_petroleo_lp_modelo10_kilian" / "graficos_ols_limpos"
DUPLA  = BASE / "output_petroleo_lp_modelo10_kilian" / "graficos_ols_dupla_faixa"
SD_DIR = BASE / "output_petroleo_lp_state_dependent_Modelo12"
OUTPUT = BASE / "TCC_Apresentacao_Final.pptx"

# ─── Paleta de Cores ────────────────────────────────────────────────────────
BG        = RGBColor(0xF7, 0xF9, 0xFC)   # fundo principal (branco-azulado)
NAVY      = RGBColor(0x0D, 0x21, 0x37)   # azul marinho escuro (titulos)
TEAL      = RGBColor(0x1A, 0x6B, 0x8A)   # azul-petroleo (destaques)
TEAL_LT   = RGBColor(0xD0, 0xEA, 0xF2)   # teal claro (fundo de caixas)
RED       = RGBColor(0xC0, 0x39, 0x2B)   # vermelho (destaque critico)
RED_LT    = RGBColor(0xFA, 0xE5, 0xE3)   # vermelho claro (fundo de alerta)
GOLD      = RGBColor(0xD4, 0xA0, 0x17)   # dourado (metricas)
GOLD_LT   = RGBColor(0xFD, 0xF3, 0xD8)   # dourado claro
GREEN     = RGBColor(0x1E, 0x8A, 0x44)   # verde (confirmado)
GREEN_LT  = RGBColor(0xD5, 0xF0, 0xDF)   # verde claro
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)   # texto principal
TEXT_MID  = RGBColor(0x4A, 0x55, 0x68)   # texto secundario
TEXT_LITE = RGBColor(0x71, 0x80, 0x96)   # texto terciario
DIVIDER   = RGBColor(0xCF, 0xD8, 0xE3)   # linha divisoria

# Slide: 33.867 cm x 19.05 cm  =  13.333" x 7.5"
SW = Inches(13.333)
SH = Inches(7.5)


# ─── Helpers ────────────────────────────────────────────────────────────────

def new_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl

def rect(sl, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def txt(sl, text, l, t, w, h,
        size=Pt(14), bold=False, italic=False,
        color=TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def txt2(sl, line1, line2, l, t, w, h,
         s1=Pt(14), s2=Pt(11), b1=True, b2=False,
         c1=TEXT_DARK, c2=TEXT_MID, align=PP_ALIGN.LEFT):
    """Caixa de texto com dois paragrafos (titulo + descricao)."""
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = line1
    r.font.name = "Calibri"; r.font.size = s1; r.font.bold = b1
    r.font.color.rgb = c1

    p2 = tf.add_paragraph()
    p2.alignment = align
    r2 = p2.add_run()
    r2.text = line2
    r2.font.name = "Calibri"; r2.font.size = s2; r2.font.bold = b2
    r2.font.color.rgb = c2

def pic(sl, path, l, t, w, h):
    p = Path(path)
    if not p.exists():
        print(f"  [NAO ENCONTRADO] {p.name}"); return
    sl.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))
    print(f"  + {p.name}")

def header_bar(sl, title, subtitle=None):
    """Barra de cabecalho padrao: linha teal + titulo + subtitulo."""
    rect(sl, 0, 0, 13.333, 0.08, fill=TEAL)
    txt(sl, title, 0.45, 0.15, 12.5, 0.7,
        size=Pt(28), bold=True, color=NAVY)
    if subtitle:
        txt(sl, subtitle, 0.45, 0.82, 12.5, 0.38,
            size=Pt(13), italic=True, color=TEXT_MID)

def footer(sl, num, total=12):
    rect(sl, 0, 7.28, 13.333, 0.01, fill=DIVIDER)
    txt(sl, "Pedro Franck Minella  |  TCC Ibmec-DF  |  2026",
        0.45, 7.3, 9, 0.22, size=Pt(8), color=TEXT_LITE)
    txt(sl, f"{num}/{total}", 12.2, 7.3, 1.0, 0.22,
        size=Pt(8), color=TEXT_LITE, align=PP_ALIGN.RIGHT)

def number_box(sl, num_str, label, l, t, w=2.8, h=1.0,
               num_color=TEAL, box_fill=TEAL_LT):
    rect(sl, l, t, w, h, fill=box_fill,
         line=TEAL, lw=Pt(1.5))
    txt(sl, num_str, l, t + 0.05, w, 0.45,
        size=Pt(26), bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(sl, label, l, t + 0.48, w, 0.48,
        size=Pt(9.5), color=TEXT_MID, align=PP_ALIGN.CENTER)

def info_box(sl, title, body, l, t, w, h,
             title_color=TEAL, body_color=TEXT_MID,
             bg=WHITE, border=TEAL, border_w=Pt(1)):
    rect(sl, l, t, w, h, fill=bg, line=border, lw=border_w)
    rect(sl, l, t, 0.05, h, fill=border)   # borda esquerda colorida
    txt(sl, title, l + 0.12, t + 0.08, w - 0.18, 0.32,
        size=Pt(11), bold=True, color=title_color)
    txt(sl, body, l + 0.12, t + 0.40, w - 0.18, h - 0.5,
        size=Pt(10), color=body_color)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ═══════════════════════════════════════════════════════════════════════════
def slide1(prs):
    sl = new_slide(prs)

    # Faixa lateral esquerda (teal escuro)
    rect(sl, 0, 0, 0.6, 7.5, fill=NAVY)
    # Linha dourada vertical
    rect(sl, 0.6, 0, 0.06, 7.5, fill=TEAL)

    # Titulo
    txt(sl, "Efeitos dos Choques do Petróleo\nsobre a Inflação no Brasil",
        0.9, 1.0, 11.8, 2.2, size=Pt(36), bold=True, color=NAVY)

    # Subtitulo
    txt(sl, "Evidências via Combustíveis, Petrobras e Local Projections (2003-2026)",
        0.9, 3.2, 11.0, 0.55, size=Pt(15), italic=True, color=TEAL)

    # Linha divisoria
    rect(sl, 0.9, 3.85, 11.8, 0.025, fill=DIVIDER)

    # Autor e orientador
    txt(sl, "Pedro Franck Minella",
        0.9, 3.95, 8, 0.5, size=Pt(18), bold=True, color=TEXT_DARK)
    txt(sl, "Orientador: Prof. Silvio Costa   |   Ibmec-DF, 2026",
        0.9, 4.45, 8, 0.38, size=Pt(12), color=TEXT_MID)

    # Metricas em boxes
    mets = [
        ("23 anos", "Amostra\n2003-2026"),
        ("4 modelos", "Especificações\neconométricas"),
        ("6 variáveis", "Preços de\ncombustíveis"),
        ("h = 12", "Horizonte de\nprojeção (meses)"),
    ]
    for i, (val, lbl) in enumerate(mets):
        bx = 0.9 + i * 3.1
        rect(sl, bx, 5.2, 2.85, 1.5, fill=WHITE, line=DIVIDER, lw=Pt(1))
        rect(sl, bx, 5.2, 2.85, 0.08, fill=TEAL)
        txt(sl, val, bx, 5.35, 2.85, 0.65,
            size=Pt(28), bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        txt(sl, lbl, bx, 5.98, 2.85, 0.65,
            size=Pt(10), color=TEXT_MID, align=PP_ALIGN.CENTER)

    footer(sl, 1)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — O ENIGMA DO REPASSE BRASILEIRO
# ═══════════════════════════════════════════════════════════════════════════
def slide2(prs):
    sl = new_slide(prs)
    header_bar(sl, "O Enigma do Repasse Brasileiro",
               "Por que o Brasil não segue a teoria macro clássica?")

    # Coluna esquerda: teoria classica
    rect(sl, 0.4, 1.3, 5.6, 5.5, fill=WHITE, line=DIVIDER, lw=Pt(1))
    rect(sl, 0.4, 1.3, 5.6, 0.06, fill=TEXT_MID)
    txt(sl, "A Teoria Macro Clássica",
        0.55, 1.4, 5.3, 0.45, size=Pt(14), bold=True, color=TEXT_DARK)

    passos_esq = [
        ("Choque Externo (Brent sobe)", TEXT_DARK),
        ("Câmbio amplifica em R$", TEXT_MID),
        ("Combustíveis encarecem", TEXT_MID),
        ("Inflação sobe imediatamente", TEXT_MID),
    ]
    setas = ["↓", "↓", "↓"]
    for i, (p_, c) in enumerate(passos_esq):
        yb = 1.95 + i * 1.12
        rect(sl, 0.6, yb, 5.2, 0.68, fill=TEAL_LT, line=TEAL, lw=Pt(1))
        txt(sl, p_, 0.72, yb + 0.12, 5.0, 0.44, size=Pt(12), color=c)
        if i < 3:
            txt(sl, "↓", 2.8, yb + 0.7, 0.6, 0.35,
                size=Pt(16), bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    txt(sl, "Repasse mecânico e automático.",
        0.55, 5.55, 5.3, 0.35, size=Pt(10), italic=True, color=TEXT_LITE)

    # Divisor central
    rect(sl, 6.3, 1.3, 0.025, 5.5, fill=DIVIDER)
    txt(sl, "vs.", 6.05, 3.8, 0.55, 0.5,
        size=Pt(16), bold=True, color=TEXT_MID, align=PP_ALIGN.CENTER)

    # Coluna direita: realidade brasileira
    rect(sl, 6.6, 1.3, 6.3, 5.5, fill=WHITE, line=RED, lw=Pt(1.5))
    rect(sl, 6.6, 1.3, 6.3, 0.06, fill=RED)
    txt(sl, "A Realidade Brasileira",
        6.75, 1.4, 6.1, 0.45, size=Pt(14), bold=True, color=RED)

    fricoes = [
        ("Câmbio (R$/USD)", "Reage a choques fiscais internos — "
         "não só ao petróleo. Pode amplificar OU atenuar."),
        ("Petrobras (Refinaria)", "Decide quando e quanto repassar. "
         "Antes de 2016: bloqueava. Após PPI: repassa imediatamente."),
        ("Tributos e Misturas", "CIDE, ICMS, mistura de etanol alteram "
         "a magnitude final do repasse na bomba."),
        ("Mercado e Transporte", "Fricções contratuais e tarifas "
         "administradas retardam a chegada ao IPCA."),
    ]
    for i, (titulo, desc) in enumerate(fricoes):
        yb = 1.95 + i * 1.2
        info_box(sl, titulo, desc, 6.75, yb, 6.0, 1.05,
                 title_color=RED, body_color=TEXT_MID, bg=RED_LT, border=RED)

    footer(sl, 2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — O MECANISMO DE TRANSMISSAO
# ═══════════════════════════════════════════════════════════════════════════
def slide3(prs):
    sl = new_slide(prs)
    header_bar(sl, "O Mecanismo de Transmissão",
               "Cadeia causal do choque ao consumidor — cada elo é medido por Local Projections")

    nos = [
        ("🛢", "Brent\n(USD)", TEAL),
        ("💱", "Câmbio\n(R$/USD)", TEAL),
        ("🏭", "Petrobras\n(Refinaria)", NAVY),
        ("⛽", "Bomba\n(Gasolina/Diesel)", TEAL),
        ("📊", "IPCA\nTransportes", GREEN),
    ]
    bw, bh = 2.1, 1.55
    gap = 0.4
    total_w = len(nos) * bw + (len(nos) - 1) * gap
    start_x = (13.333 - total_w) / 2
    y_box = 1.6

    for i, (icon, label_, color) in enumerate(nos):
        bx = start_x + i * (bw + gap)
        rect(sl, bx, y_box, bw, bh, fill=WHITE, line=color, lw=Pt(2))
        rect(sl, bx, y_box, bw, 0.07, fill=color)
        txt(sl, icon, bx, y_box + 0.1, bw, 0.55,
            size=Pt(22), align=PP_ALIGN.CENTER)
        txt(sl, label_, bx, y_box + 0.65, bw, 0.82,
            size=Pt(11), bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)
        if i < len(nos) - 1:
            ax = bx + bw + 0.05
            txt(sl, "→", ax, y_box + 0.5, gap - 0.05, 0.55,
                size=Pt(20), bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Descricoes abaixo de cada no
    descs = [
        "Choque global\nem dólares",
        "Converte e\namplifica",
        "Define timing\ne magnitude",
        "Impostos\ne misturas",
        "Peso direto\nna cesta",
    ]
    for i, d in enumerate(descs):
        bx = start_x + i * (bw + gap)
        txt(sl, d, bx, y_box + bh + 0.12, bw, 0.65,
            size=Pt(9.5), italic=True, color=TEXT_LITE, align=PP_ALIGN.CENTER)

    # Caixa de hipotese central
    rect(sl, 0.4, 4.15, 12.5, 1.2, fill=TEAL_LT, line=TEAL, lw=Pt(1.5))
    rect(sl, 0.4, 4.15, 0.07, 1.2, fill=TEAL)
    txt(sl, "Hipótese Central",
        0.6, 4.2, 4, 0.38, size=Pt(11), bold=True, color=TEAL)
    txt(sl,
        "O repasse é setorialmente concentrado no IPCA Transportes. A política de preços da "
        "Petrobras atua como válvula reguladora — bloqueando (pré-2016) ou acelerando (pós-2016 PPI) "
        "a velocidade e magnitude desta transmissão.",
        0.6, 4.55, 12.2, 0.7, size=Pt(11), color=TEXT_DARK)

    # 3 colunas de resumo
    cols = [
        ("Filtro Cambial", "Converte o choque em R$, amplia volatilidade doméstica.", TEAL),
        ("Filtro Petrobras", "O pilar central desta pesquisa. Decide se o choque passa.", NAVY),
        ("Filtro de Mercado", "Distribuição, tributos e lei de substituição etanol/gasolina.", GOLD),
    ]
    for i, (ti, de, co) in enumerate(cols):
        bx = 0.4 + i * 4.3
        info_box(sl, ti, de, bx, 5.55, 4.1, 1.55,
                 title_color=co, bg=WHITE, border=co)

    footer(sl, 3)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — METODOLOGIA + GRAFICO REAL IPCA TRANSPORTES
# ═══════════════════════════════════════════════════════════════════════════
def slide4(prs):
    sl = new_slide(prs)
    header_bar(sl, "A Lente Metodológica: Local Projections (Jordà, 2005)",
               "Estimação direta de y(t+h) = α(h) + β(h)·ΔBrent(t) + controles + ε(t+h) para h = 0,...,12")

    # Coluna esquerda: bullets metodologicos
    rect(sl, 0.4, 1.3, 5.3, 5.9, fill=WHITE, line=DIVIDER, lw=Pt(1))
    rect(sl, 0.4, 1.3, 5.3, 0.06, fill=TEAL)

    pontos = [
        ("Por que LP e não VAR?",
         "O VAR itera erros ao longo do tempo. O LP estima cada horizonte h de forma independente "
         "por OLS, tornando-o robusto a erros de especificação."),
        ("Correção HAC (Newey-West)",
         "LP acumula autocorrelação MA(h) por construção. Correção com maxlags=h garante "
         "erros-padrão corretos e bandas de confiança precisas."),
        ("Controles Incluídos",
         "Câmbio R$/USD, IBC-Br (atividade), Selic, Expectativas Focus e Índice Kilian "
         "(demanda global de commodities)."),
        ("Horizonte",
         "h = 0 a 12 meses. Cada regressão é independente — sem propagação de erros "
         "entre horizontes."),
    ]
    for i, (ti, de) in enumerate(pontos):
        yb = 1.45 + i * 1.35
        info_box(sl, ti, de, 0.5, yb, 5.1, 1.22,
                 title_color=TEAL, bg=TEAL_LT, border=TEAL)

    # Coluna direita: grafico real IPCA Transportes
    rect(sl, 6.0, 1.3, 6.9, 5.9, fill=WHITE, line=DIVIDER, lw=Pt(1))
    rect(sl, 6.0, 1.3, 6.9, 0.06, fill=TEAL)
    txt(sl, "Exemplo: Gasolina C → IPCA Transportes (h=0 a 12)",
        6.05, 1.32, 6.8, 0.38,
        size=Pt(10), bold=True, color=TEAL)
    pic(sl, DUPLA / "LP_OLS_Kilian_dupla_faixa_ipca_transporte.png",
        6.05, 1.72, 6.8, 5.35)

    # Nota de fonte
    txt(sl, "Fonte: Elaboracao propria. Controles: cambio, IBC-Br, Selic, Focus, Kilian. IC 90% e 95%.",
        6.05, 7.12, 6.8, 0.25,
        size=Pt(7.5), italic=True, color=TEXT_LITE)

    footer(sl, 4)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PASSO 1: A CHEGADA À BOMBA (3 graficos reais)
# ═══════════════════════════════════════════════════════════════════════════
def slide5(prs):
    sl = new_slide(prs)
    header_bar(sl, "Passo 1: A Chegada à Bomba",
               "O canal primário está ativo — choque do petróleo se transmite aos combustíveis com ritmos distintos")

    # 3 graficos em linha horizontal
    graficos = [
        ("Diesel", DUPLA / "LP_OLS_Kilian_dupla_faixa_diesel.png",
         "Repasse ultra-rápido (h=0,1), absorvido em ~3 meses.", GOLD),
        ("Gasolina A — Refinaria", DUPLA / "LP_OLS_Kilian_dupla_faixa_gasolina_refinaria.png",
         "Acumulação gradual em platôs. Reflete janelas decisórias da Petrobras.", TEAL),
        ("Gasolina C — Bomba", DUPLA / "LP_OLS_Kilian_dupla_faixa_gasolina.png",
         "Impacto violento, imediato (h=0) e duradouro até h=12. *** em todos os horizontes.", GREEN),
    ]

    gw = 4.0
    gh = 4.9
    for i, (titulo, imgpath, nota, cor) in enumerate(graficos):
        bx = 0.4 + i * (gw + 0.2)
        # Container
        rect(sl, bx, 1.3, gw, gh + 0.9, fill=WHITE, line=DIVIDER, lw=Pt(1))
        rect(sl, bx, 1.3, gw, 0.06, fill=cor)
        # Titulo do grafico
        txt(sl, titulo, bx + 0.08, 1.32, gw - 0.15, 0.38,
            size=Pt(11), bold=True, color=cor)
        # Imagem
        pic(sl, imgpath, bx + 0.05, 1.72, gw - 0.1, gh)
        # Nota abaixo
        txt(sl, nota, bx + 0.08, 6.72, gw - 0.15, 0.45,
            size=Pt(8.5), italic=True, color=TEXT_MID)

    # Caixa de conclusao base
    rect(sl, 0.4, 7.0, 12.5, 0.38, fill=GREEN_LT, line=GREEN, lw=Pt(1))
    txt(sl,
        "Confirmado: Todos os combustíveis respondem positiva e significativamente ao choque do Brent. "
        "Gasolina C apresenta o maior e mais persistente pass-through (***  em h=0 a h=12).",
        0.55, 7.02, 12.2, 0.34,
        size=Pt(9.5), bold=True, color=GREEN)

    footer(sl, 5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PASSO 2: O EPICENTRO (IPCA TRANSPORTES com 43%)
# ═══════════════════════════════════════════════════════════════════════════
def slide6(prs):
    sl = new_slide(prs)
    header_bar(sl, "Passo 2: O Epicentro — IPCA Transportes",
               "A inflação setorial: resposta acumulada da Gasolina C sobre o IPCA Transportes")

    # Grafico grande ocupando quase o slide inteiro
    rect(sl, 0.4, 1.3, 9.1, 5.9, fill=WHITE, line=DIVIDER, lw=Pt(1))
    rect(sl, 0.4, 1.3, 9.1, 0.06, fill=RED)
    txt(sl, "Gasolina C → IPCA Transportes  |  IC 90% e 95%  |  h=0 a 12 meses",
        0.5, 1.32, 9.0, 0.38, size=Pt(10), bold=True, color=RED)
    pic(sl, DUPLA / "LP_OLS_Kilian_dupla_faixa_ipca_transporte.png",
        0.4, 1.7, 9.1, 5.35)
    txt(sl, "Fonte: LP-OLS com Controle Kilian. HAC Newey-West. Transformacoes em ln x 100.",
        0.45, 7.1, 9.0, 0.25, size=Pt(7.5), italic=True, color=TEXT_LITE)

    # Painel lateral de metricas
    metricas = [
        ("0,269", "Pass-through\nih=0 (imediato)", RED),
        ("0,431", "Pass-through\nih=12 (persistente)", RED),
        ("~43%", "do choque\ntransferido", TEAL),
        ("t > 3,6", "t-stat em\ntodo horizonte", GREEN),
    ]
    for i, (val, lbl, cor) in enumerate(metricas):
        yb = 1.3 + i * 1.52
        rect(sl, 9.75, yb, 3.1, 1.38, fill=WHITE, line=cor, lw=Pt(1.5))
        rect(sl, 9.75, yb, 3.1, 0.06, fill=cor)
        txt(sl, val, 9.75, yb + 0.1, 3.1, 0.72,
            size=Pt(32), bold=True, color=cor, align=PP_ALIGN.CENTER)
        txt(sl, lbl, 9.75, yb + 0.82, 3.1, 0.5,
            size=Pt(10), color=TEXT_MID, align=PP_ALIGN.CENTER)

    footer(sl, 6)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — A DILUICAO MACRO (3 graficos: Transportes esq., Geral e Diesel dir.)
# ═══════════════════════════════════════════════════════════════════════════
def slide7(prs):
    sl = new_slide(prs)
    header_bar(sl, "A Diluição Macro: IPCA Transportes vs. IPCA Geral",
               "O choque gera inflação aguda no setor, mas se dilui na vasta cesta de consumo do brasileiro")

    # Grafico esquerdo grande: IPCA Transportes
    rect(sl, 0.4, 1.3, 6.2, 5.9, fill=WHITE, line=RED, lw=Pt(1.5))
    rect(sl, 0.4, 1.3, 6.2, 0.06, fill=RED)
    txt(sl, "IPCA Transportes  ✓  (pass-through ~43%)",
        0.5, 1.32, 6.1, 0.38, size=Pt(10), bold=True, color=RED)
    pic(sl, LIMPOS / "LP_OLS_Kilian_limpo_ipca_transporte.png",
        0.4, 1.7, 6.2, 5.35)

    # Grafico direito superior: IPCA Geral
    rect(sl, 6.9, 1.3, 6.0, 2.85, fill=WHITE, line=TEXT_MID, lw=Pt(1))
    rect(sl, 6.9, 1.3, 6.0, 0.06, fill=TEXT_MID)
    txt(sl, "IPCA Geral  (efeito diluido — p<0,05 apenas em h=0 a h=4)",
        7.0, 1.32, 5.9, 0.38, size=Pt(10), bold=True, color=TEXT_MID)
    pic(sl, LIMPOS / "LP_OLS_Kilian_limpo_ipca_geral.png",
        6.9, 1.7, 6.0, 2.35)

    # Grafico direito inferior: Diesel IPCA geral
    rect(sl, 6.9, 4.3, 6.0, 2.85, fill=WHITE, line=GOLD, lw=Pt(1))
    rect(sl, 6.9, 4.3, 6.0, 0.06, fill=GOLD)
    txt(sl, "Diesel → IPCA Transportes  (choque temporario, absorvido em ~3 meses)",
        7.0, 4.32, 5.9, 0.38, size=Pt(10), bold=True, color=GOLD)
    pic(sl, LIMPOS / "LP_OLS_Kilian_limpo_diesel.png",
        6.9, 4.7, 6.0, 2.35)

    # Caixa de conclusao
    rect(sl, 0.4, 7.22, 12.5, 0.38, fill=TEAL_LT, line=TEAL, lw=Pt(1))
    txt(sl,
        "Gasolina C gera impacto 7x maior no IPCA Transportes que no IPCA Geral. "
        "IPCA Geral perde significância após o 5.º mês — choque de custo setorial transitório.",
        0.55, 7.24, 12.2, 0.34,
        size=Pt(9.5), bold=True, color=TEAL)

    footer(sl, 7)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — A HIPOTESE INGENUA vs. ARQUITETURA REAL DA TRANSMISSAO
# ═══════════════════════════════════════════════════════════════════════════
def slide8(prs):
    sl = new_slide(prs)
    header_bar(sl, "A Hipótese Ingênua (Incorreta) vs. Arquitetura Real",
               "O modelo do TCC revela: o petróleo não causa inflação diretamente — o canal passa por múltiplos filtros")

    # ── PARTE SUPERIOR: Hipótese Ingênua ─────────────────────────────────
    rect(sl, 0.4, 1.28, 12.5, 1.85, fill=RED_LT, line=RED, lw=Pt(1.5))
    rect(sl, 0.4, 1.28, 12.5, 0.06, fill=RED)
    txt(sl, "A Hipótese Ingênua (INCORRETA)",
        0.55, 1.3, 6, 0.38, size=Pt(13), bold=True, color=RED)

    # Caixa Petróleo
    rect(sl, 0.7, 1.75, 2.5, 0.85, fill=WHITE, line=RED, lw=Pt(1.5))
    txt(sl, "Petróleo (USD)", 0.75, 1.88, 2.4, 0.6,
        size=Pt(12), bold=True, color=RED, align=PP_ALIGN.CENTER)

    # Cruz vermelha (X) no centro
    txt(sl, "✕", 3.5, 1.77, 2.5, 0.85,
        size=Pt(36), bold=True, color=RED, align=PP_ALIGN.CENTER)
    txt(sl, "SIMPLIFICAÇÃO INCORRETA",
        3.55, 2.38, 2.5, 0.35,
        size=Pt(7.5), italic=True, color=RED, align=PP_ALIGN.CENTER)

    # Seta -> Inflação Geral
    txt(sl, "→", 6.1, 1.77, 1.0, 0.85,
        size=Pt(28), bold=True, color=RED, align=PP_ALIGN.CENTER)
    rect(sl, 7.2, 1.75, 3.0, 0.85, fill=WHITE, line=RED, lw=Pt(1.5))
    txt(sl, "Inflação Geral", 7.25, 1.88, 2.9, 0.6,
        size=Pt(12), bold=True, color=RED, align=PP_ALIGN.CENTER)

    txt(sl, "A transmissão NÃO é mecânica nem automática.",
        10.4, 1.75, 2.3, 0.85,
        size=Pt(9.5), italic=True, color=RED)

    # ── PARTE INFERIOR: Arquitetura Real ─────────────────────────────────
    rect(sl, 0.4, 3.28, 12.5, 3.85, fill=TEAL_LT, line=TEAL, lw=Pt(1.5))
    rect(sl, 0.4, 3.28, 12.5, 0.06, fill=TEAL)
    txt(sl, "Arquitetura Real da Transmissão (Modelo do TCC)",
        0.55, 3.3, 10, 0.38, size=Pt(13), bold=True, color=TEAL)

    # Nos da cadeia
    nos_real = [
        ("Node 1\nOrigem",       "Brent (USD)\n+ Choque\nde Oferta",   NAVY,  0.55),
        ("Node 2\nFiltro 1",     "Câmbio\n(BRL)",                       TEAL,  2.90),
        ("Node 3\nVálvula",      "Refinaria\nPetrobras\n(Política)",    RED,   5.25),
        ("Node 4\nRamificação",  "Diesel\nGasolina\nEtanol",            GOLD,  7.60),
        ("Node 5\nImpacto Final","IPCA\nTransportes\n(Concentrado)",   GREEN, 10.15),
    ]
    for header, body, cor, bx in nos_real:
        rect(sl, bx, 3.75, 2.1, 1.95, fill=WHITE, line=cor, lw=Pt(1.5))
        rect(sl, bx, 3.75, 2.1, 0.06, fill=cor)
        txt(sl, header, bx + 0.05, 3.77, 2.0, 0.35,
            size=Pt(7.5), bold=True, color=cor, align=PP_ALIGN.CENTER)
        txt(sl, body, bx + 0.05, 4.12, 2.0, 1.5,
            size=Pt(11), bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # Setas entre nos
    for bx in [2.65, 5.0, 7.35, 9.9]:
        txt(sl, "→", bx, 4.35, 0.55, 0.55,
            size=Pt(20), bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Resultado inferior: IPCA Geral diluido
    rect(sl, 10.15, 5.82, 2.1, 0.88, fill=WHITE, line=TEXT_LITE, lw=Pt(1))
    txt(sl, "IPCA Geral\n(Efeito Diluído)",
        10.2, 5.9, 2.0, 0.72,
        size=Pt(10), color=TEXT_MID, align=PP_ALIGN.CENTER)
    txt(sl, "↓", 11.15, 5.68, 0.45, 0.35,
        size=Pt(16), bold=True, color=TEXT_LITE, align=PP_ALIGN.CENTER)

    # Caixa de conclusao
    rect(sl, 0.4, 7.15, 12.5, 0.45, fill=NAVY)
    txt(sl,
        "O desafio econométrico é isolar o timing e a magnitude do repasse etapa a etapa, "
        "controlando ruídos macroeconômicos, substituição etanol/gasolina e quebras institucionais.",
        0.55, 7.2, 12.2, 0.35,
        size=Pt(10), bold=True, color=WHITE)

    footer(sl, 8, total=13)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8B — TABELA DE SIGNIFICANCIA h=0 a h=6
# ═══════════════════════════════════════════════════════════════════════════
def slide8b(prs):
    sl = new_slide(prs)
    header_bar(sl, "Resultados: Coeficientes e Significância (h=0 a h=6)",
               "Modelo LP-OLS com Controle Kilian (Modelo 10). Erros-padrão HAC Newey-West. *** p<0,01  ** p<0,05  * p<0,10")

    def sig(p):
        if p < 0.01:  return "***"
        if p < 0.05:  return "**"
        if p < 0.10:  return "*"
        return ""

    def fmt(v): return f"{v:.3f}"

    # Dados reais dos CSVs (coef, pvalor) para h=0..6
    # Colunas: variavel dependente | h=0 | h=1 | h=2 | h=3 | h=4 | h=5 | h=6
    dados = [
        # (rotulo, [(coef, pvalor) para h=0 a h=6])
        ("Gasolina C (Bomba)", [
            (0.6833, 6.06e-06), (1.4059, 1.05e-06), (1.6022, 3.36e-04),
            (1.5335, 1.53e-03), (1.3518, 8.45e-03), (1.2539, 1.46e-02), (1.1253, 2.46e-02),
        ]),
        ("Gasolina A (Refinaria)", [
            (0.7060, 6.58e-05), (1.4515, 9.56e-05), (1.7012, 1.47e-03),
            (1.5913, 7.85e-03), (1.4130, 3.35e-02), (1.2957, 6.68e-02), (1.1498, 9.21e-02),
        ]),
        ("Diesel", [
            (0.6995, 2.62e-06), (1.3967, 9.05e-08), (1.9018, 2.01e-06),
            (1.8608, 7.79e-05), (1.9264, 9.61e-05), (1.8987, 6.84e-04), (1.9928, 1.27e-03),
        ]),
        ("IPCA Transportes", [
            (0.1657, 9.05e-04), (0.3731, 1.40e-04), (0.4575, 1.96e-03),
            (0.4612, 4.11e-03), (0.4387, 1.23e-02), (0.4284, 2.01e-02), (0.3655, 5.48e-02),
        ]),
        ("IPCA Geral", [
            (-0.0038, 0.7815), (0.0399, 0.1467), (0.0541, 0.2454),
            (0.0627, 0.2695), (0.0583, 0.3403), (0.0520, 0.4396), (0.0153, 0.8270),
        ]),
    ]

    # Cores por variavel
    cores = [GREEN, TEAL, GOLD, RED, TEXT_MID]

    # Dimensoes da tabela
    col_w  = [3.05, 1.45, 1.45, 1.45, 1.45, 1.45, 1.45, 1.45]  # [variavel, h0..h6]
    row_h  = 0.72
    x0, y0 = 0.35, 1.3

    # Cabecalho da tabela
    headers = ["Variável Dependente", "h = 0", "h = 1", "h = 2", "h = 3", "h = 4", "h = 5", "h = 6"]
    cx = x0
    for j, (hdr, cw) in enumerate(zip(headers, col_w)):
        rect(sl, cx, y0, cw - 0.02, row_h * 0.7, fill=NAVY)
        txt(sl, hdr, cx + 0.05, y0 + 0.08, cw - 0.1, row_h * 0.55,
            size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += cw

    # Linhas de dados
    for i, (var_lbl, vals) in enumerate(dados):
        cor = cores[i]
        yb = y0 + row_h * 0.7 + i * row_h
        bg_row = TEAL_LT if i % 2 == 0 else WHITE

        # Celula do nome da variavel
        rect(sl, x0, yb, col_w[0] - 0.02, row_h - 0.03,
             fill=bg_row, line=cor, lw=Pt(1.5))
        rect(sl, x0, yb, 0.06, row_h - 0.03, fill=cor)  # borda esq colorida
        txt(sl, var_lbl, x0 + 0.1, yb + 0.12, col_w[0] - 0.15, row_h - 0.1,
            size=Pt(11), bold=True, color=cor)

        # Celulas de coeficientes
        cx = x0 + col_w[0]
        for j, (coef, pval) in enumerate(vals):
            s = sig(pval)
            c_bg = GREEN_LT if pval < 0.01 else (GOLD_LT if pval < 0.05 else (RED_LT if pval < 0.10 else WHITE))
            rect(sl, cx, yb, col_w[j+1] - 0.02, row_h - 0.03,
                 fill=c_bg, line=DIVIDER, lw=Pt(0.5))
            # Coeficiente
            txt(sl, fmt(coef), cx + 0.02, yb + 0.05, col_w[j+1] - 0.05, 0.35,
                size=Pt(11), bold=(pval < 0.05), color=TEXT_DARK, align=PP_ALIGN.CENTER)
            # Estrelas
            txt(sl, s, cx + 0.02, yb + 0.38, col_w[j+1] - 0.05, 0.28,
                size=Pt(10), bold=True,
                color=GREEN if pval<0.01 else (GOLD if pval<0.05 else (RED if pval<0.10 else TEXT_LITE)),
                align=PP_ALIGN.CENTER)
            cx += col_w[j+1]

    # Legenda de cores
    y_leg = y0 + row_h * 0.7 + len(dados) * row_h + 0.12
    legenda = [
        (GREEN_LT, GREEN, "p < 0,01 (***)"),
        (GOLD_LT,  GOLD,  "p < 0,05 (**) "),
        (RED_LT,   RED,   "p < 0,10 (*)  "),
        (WHITE,    DIVIDER, "ns            "),
    ]
    for k, (bg, brd, lbl) in enumerate(legenda):
        lx = 0.35 + k * 3.1
        rect(sl, lx, y_leg, 0.45, 0.35, fill=bg, line=brd, lw=Pt(1))
        txt(sl, lbl, lx + 0.5, y_leg + 0.05, 2.5, 0.3,
            size=Pt(9.5), color=TEXT_MID)

    txt(sl, "Nota: Coeficientes representam p.p. acumulados de resposta do preco do combustivel (ou do IPCA) "
            "a um aumento de 1 p.p. no preco da Gasolina C. Transformacoes: dln x 100.",
        0.35, y_leg + 0.45, 12.6, 0.38,
        size=Pt(8), italic=True, color=TEXT_LITE)

    footer(sl, 9, total=13)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — O CLIMAX: A VALVULA DA PETROBRAS (grafico real SD_LP)
# ═══════════════════════════════════════════════════════════════════════════
def slide9(prs):
    sl = new_slide(prs)
    header_bar(sl, "O Clímax: A Válvula da Petrobras",
               "A quebra institucional de setembro/2016 alterou a física do repasse inflacionário (Teste de Wald p<0,05)")

    # Grafico principal (SD_LP — lado esquerdo)
    rect(sl, 0.4, 1.3, 7.6, 5.9, fill=WHITE, line=NAVY, lw=Pt(1.5))
    rect(sl, 0.4, 1.3, 7.6, 0.06, fill=NAVY)
    txt(sl, "IPCA Transportes: Pré-2016 (azul) vs. Pós-2016/PPI (vermelho)",
        0.5, 1.32, 7.5, 0.38, size=Pt(10), bold=True, color=NAVY)
    pic(sl, SD_DIR / "SD_LP_ipca_transporte_mensal.png",
        0.4, 1.7, 7.6, 5.35)
    txt(sl, "Fonte: Modelo LP State-Dependent (Modelo 12). Corte: set/2016 (PPI Petrobras). IC 90%.",
        0.45, 7.1, 7.5, 0.25, size=Pt(7.5), italic=True, color=TEXT_LITE)

    # Painel direito: os dois Brasis
    txt(sl, "Dois Brasis Institucionais:",
        8.3, 1.3, 4.6, 0.42, size=Pt(14), bold=True, color=NAVY)

    info_box(sl, "Fase Azul — Pré-2016 (Precos Administrados)",
             "A Petrobras fecha a válvula. O choque internacional morre nos "
             "portões da refinaria. Repasse estatisticamente nulo em todos os horizontes.",
             8.3, 1.82, 4.6, 1.55,
             title_color=TEAL, bg=TEAL_LT, border=TEAL)

    info_box(sl, "Fase Vermelha — Pós-2016 (PPI)",
             "Adoção da Política de Paridade de Importação. O choque fura o "
             "bloqueio com fluxo livre. Repasse imediato e robusto em h=0 a h=12.",
             8.3, 3.52, 4.6, 1.55,
             title_color=RED, bg=RED_LT, border=RED)

    # Box Wald
    rect(sl, 8.3, 5.22, 4.6, 1.25, fill=NAVY, line=NAVY)
    txt(sl, "Teste de Wald",
        8.45, 5.3, 4.3, 0.38, size=Pt(11), bold=True, color=WHITE)
    txt(sl, "p < 0,05", 8.45, 5.65, 2.0, 0.55,
        size=Pt(28), bold=True, color=GOLD)
    txt(sl, "Diferença estrutural\nentre regimes confirmada.",
        10.35, 5.65, 2.4, 0.7, size=Pt(10), color=WHITE)

    # Implicacao
    rect(sl, 8.3, 6.6, 4.6, 0.75, fill=RED_LT, line=RED, lw=Pt(1))
    txt(sl,
        "Após 2016, o Brasil tornou-se mais vulnerável "
        "a tensões geopolíticas globais que antes ficavam retidas na Petrobras.",
        8.42, 6.66, 4.4, 0.62, size=Pt(9.5), color=TEXT_DARK)

    footer(sl, 10, total=13)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — O TESTE DE ESTRESSE (ROBUSTEZ)
# ═══════════════════════════════════════════════════════════════════════════
def slide10(prs):
    sl = new_slide(prs)
    header_bar(sl, "O Teste de Estresse: Robustez",
               "Validando a causalidade e a dinâmica do sistema com duas estratégias independentes")

    # Dois paineis lado a lado
    paineis = [
        {
            "titulo": "Abordagem VAR (Robustez Dinamica)",
            "cor": TEAL,
            "itens": [
                ("Sistema autorregressivo confirma integralmente a hierarquia das "
                 "respostas das Local Projections.", True),
                ("IPCA Transportes reage forte; IPCA Geral permanece "
                 "estatisticamente insignificante após h=5.", True),
                ("Ordenação de Cholesky: Brent → Câmbio → Combustível → IPCA.", False),
                ("Resultados do VAR disponíveis no Apêndice A do TCC.", False),
            ]
        },
        {
            "titulo": "Filtro de Endogeneidade (LP-IV)",
            "cor": RED,
            "itens": [
                ("Instrumento: Oil Supply News Shock (Känzig, 2021) — surpresas "
                 "de oferta da OPEP, ortogonais à demanda brasileira.", True),
                ("Estatística F do 1.º Estágio: ~120  >>>  limiar de instrumento "
                 "fraco de Staiger & Stock (F > 10).", True),
                ("Coeficientes LP-IV ≈ LP-OLS em todos os horizontes — "
                 "viés de endogeneidade é desprezível.", True),
                ("Conclusão: relação medida é causal, robusta e resistente.", False),
            ]
        }
    ]

    for j, painel in enumerate(paineis):
        bx = 0.4 + j * 6.5
        cor = painel["cor"]
        rect(sl, bx, 1.3, 6.1, 5.9, fill=WHITE, line=cor, lw=Pt(1.5))
        rect(sl, bx, 1.3, 6.1, 0.06, fill=cor)
        txt(sl, painel["titulo"], bx + 0.1, 1.32, 5.9, 0.38,
            size=Pt(12), bold=True, color=cor)
        for i, (item_txt, destaque) in enumerate(painel["itens"]):
            yb = 1.85 + i * 1.28
            bg_it = TEAL_LT if (destaque and cor == TEAL) else (RED_LT if (destaque and cor == RED) else WHITE)
            rect(sl, bx + 0.1, yb, 5.9, 1.15, fill=bg_it,
                 line=cor if destaque else DIVIDER, lw=Pt(1) if destaque else Pt(0.5))
            txt(sl, "✓  " + item_txt if destaque else item_txt,
                bx + 0.22, yb + 0.08, 5.6, 0.98,
                size=Pt(10.5), bold=destaque, color=cor if destaque else TEXT_MID)

    # Box F-stat destaque
    rect(sl, 4.6, 3.75, 2.1, 1.8, fill=RED, line=RED)
    txt(sl, "F-stat", 4.6, 3.85, 2.1, 0.38,
        size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "~120", 4.6, 4.2, 2.1, 0.72,
        size=Pt(36), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(sl, "Instrumento\nFORTE", 4.6, 4.92, 2.1, 0.55,
        size=Pt(10), color=WHITE, align=PP_ALIGN.CENTER)

    # Caixa de conclusao economica
    rect(sl, 0.4, 7.22, 12.5, 0.38, fill=NAVY, line=NAVY)
    txt(sl,
        "Conclusao Economica: A relacao medida e causal, robusta e resistente a endogeneidade. "
        "LP-OLS e LP-IV produzem resultados virtualmente identicos.",
        0.55, 7.24, 12.2, 0.34,
        size=Pt(9.5), bold=True, color=WHITE)

    footer(sl, 11, total=13)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — A HISTORIA COMPLETA (3 conclusoes)
# ═══════════════════════════════════════════════════════════════════════════
def slide11(prs):
    sl = new_slide(prs)
    header_bar(sl, "A História Completa: Três Achados Centrais",
               "Resultados robustos, consistentes entre si e com a literatura macroeconômica")

    conclusoes = [
        {
            "num": "01",
            "titulo": "O Duto Está Ativo",
            "desc": ("O choque externo do Brent transborda de fato para a bomba de gasolina. "
                     "Diesel e Gasolina A respondem com força (*** em h=0 a h=2). "
                     "Gasolina C mantém significância estatística em todo o horizonte de 12 meses."),
            "stat": "*** todos os\ncombustíveis",
            "cor": TEAL,
            "bg": TEAL_LT,
        },
        {
            "num": "02",
            "titulo": "A Inflação é Setorial",
            "desc": ("O repasse esmaga o IPCA Transportes (~43% acumulado, t>3,6 em h=0 a 12), "
                     "mas se dilui e perde força estatística no IPCA Geral após o 5.º mês. "
                     "Típico de choque de custo energético setorial e transitório."),
            "stat": "Pass-through\n~43%",
            "cor": RED,
            "bg": RED_LT,
        },
        {
            "num": "03",
            "titulo": "A Instituição Importa",
            "desc": ("O PPI de 2016 foi uma quebra estrutural comprovada pelo Teste de Wald (p<0,05). "
                     "Pré-2016: repasse bloqueado. Pós-2016: imediato e robusto. "
                     "A política de preços da Petrobras é a verdadeira válvula da inflação energética no Brasil."),
            "stat": "Wald\np < 0,05",
            "cor": NAVY,
            "bg": TEAL_LT,
        },
    ]

    for i, c in enumerate(conclusoes):
        yb = 1.35 + i * 1.95
        rect(sl, 0.4, yb, 12.0, 1.78, fill=c["bg"], line=c["cor"], lw=Pt(1.5))
        rect(sl, 0.4, yb, 0.08, 1.78, fill=c["cor"])
        # Numero
        rect(sl, 0.52, yb + 0.22, 1.0, 1.2, fill=c["cor"])
        txt(sl, c["num"], 0.52, yb + 0.38, 1.0, 0.72,
            size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Titulo e desc
        txt(sl, c["titulo"], 1.65, yb + 0.1, 8.7, 0.42,
            size=Pt(13), bold=True, color=c["cor"])
        txt(sl, c["desc"], 1.65, yb + 0.52, 8.7, 1.15,
            size=Pt(10.5), color=TEXT_DARK)
        # Stat box
        rect(sl, 10.55, yb + 0.18, 1.75, 1.42, fill=c["cor"])
        txt(sl, c["stat"], 10.55, yb + 0.35, 1.75, 1.0,
            size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    footer(sl, 12, total=13)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — MUITO OBRIGADO / ENCERRAMENTO
# ═══════════════════════════════════════════════════════════════════════════
def slide12(prs):
    sl = new_slide(prs)

    # Fundo dividido: NAVY em cima, BG em baixo
    rect(sl, 0, 0, 13.333, 3.8, fill=NAVY)

    txt(sl, "Muito Obrigado.",
        0.8, 0.6, 11.8, 1.6, size=Pt(54), bold=True, color=WHITE)
    txt(sl, "Pedro Franck Minella",
        0.8, 2.2, 11.8, 0.65, size=Pt(20), color=TEAL)
    txt(sl, "Aberto para arguição da banca.",
        0.8, 2.85, 11.8, 0.55, size=Pt(14), italic=True, color=DIVIDER)

    # Linha divisoria
    rect(sl, 0.8, 4.0, 11.7, 0.04, fill=TEAL)

    # Referencias dos modelos utilizados
    txt(sl, "Modelos Utilizados",
        0.8, 4.15, 5, 0.4, size=Pt(12), bold=True, color=NAVY)
    modelos = [
        "Modelo 10 — LP-OLS com Controle Kilian (baseline principal)",
        "Modelo 12 — LP State-Dependent por Regime de Preco da Petrobras",
        "Modelo 9  — LP-IV com Oil Supply News Shock (Kanzig, 2021)",
        "Apendice  — VAR (Robustez Dinamica)",
    ]
    for i, m in enumerate(modelos):
        txt(sl, "·  " + m, 0.8, 4.62 + i * 0.47, 6.5, 0.42,
            size=Pt(10), color=TEXT_MID)

    # Referencias bibliograficas
    txt(sl, "Referencias Principais",
        7.5, 4.15, 5.4, 0.4, size=Pt(12), bold=True, color=NAVY)
    refs = [
        "Jorda, O. (2005). Estimation and Inference of IRFs by LP. AER.",
        "Kilian, L. (2009). Not All Oil Price Shocks Are Alike. AER.",
        "Kanzig, D. (2021). The Macroeconomic Effects of Oil Supply News. AER.",
        "Ramey, V. (2016). Macroeconomic Shocks and Their Propagation. HBM.",
    ]
    for i, r in enumerate(refs):
        txt(sl, r, 7.5, 4.62 + i * 0.47, 5.5, 0.42,
            size=Pt(9.5), italic=True, color=TEXT_MID)

    # Contato
    rect(sl, 0.8, 6.85, 11.7, 0.52, fill=TEAL_LT, line=TEAL, lw=Pt(1))
    txt(sl,
        "pedrofrankminella@gmail.com   |   Ibmec-DF, 2026   |   Orientador: Prof. Silvio Costa",
        0.95, 6.93, 11.4, 0.35, size=Pt(10), color=TEAL, align=PP_ALIGN.CENTER)

    footer(sl, 13, total=13)


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Gerando slides...")
    slide1(prs);  print("  [1] Capa")
    slide2(prs);  print("  [2] Enigma do Repasse")
    slide3(prs);  print("  [3] Mecanismo de Transmissao")
    slide4(prs);  print("  [4] Metodologia + grafico LP")
    slide5(prs);  print("  [5] Chegada a Bomba (3 graficos)")
    slide6(prs);  print("  [6] Epicentro IPCA Transportes")
    slide7(prs);  print("  [7] Diluicao Macro (3 graficos)")
    slide8(prs);  print("  [8] Hipotese Ingenua vs. Arquitetura Real")
    slide8b(prs); print("  [9] Tabela Significancia h=0 a h=6")
    slide9(prs);  print("  [10] Valvula da Petrobras + SD_LP")
    slide10(prs); print("  [11] Teste de Estresse")
    slide11(prs); print("  [12] Historia Completa")
    slide12(prs); print("  [13] Encerramento")

    prs.save(str(OUTPUT))
    print(f"\nSalvo em: {OUTPUT}")

if __name__ == "__main__":
    main()
