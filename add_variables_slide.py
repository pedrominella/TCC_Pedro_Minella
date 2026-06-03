# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Cores
BG        = RGBColor(0xF7, 0xF9, 0xFC)
NAVY      = RGBColor(0x0D, 0x21, 0x37)
TEAL      = RGBColor(0x1A, 0x6B, 0x8A)
TEAL_LT   = RGBColor(0xD0, 0xEA, 0xF2)
RED       = RGBColor(0xC0, 0x39, 0x2B)
RED_LT    = RGBColor(0xFA, 0xE5, 0xE3)
GOLD      = RGBColor(0xD4, 0xA0, 0x17)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_MID  = RGBColor(0x4A, 0x55, 0x68)
TEXT_LITE = RGBColor(0x71, 0x80, 0x96)
DIVIDER   = RGBColor(0xCF, 0xD8, 0xE3)
GREEN     = RGBColor(0x1E, 0x8A, 0x44)

def rect(sl, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def txt(sl, text, l, t, w, h, size=Pt(14), bold=False, italic=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def header_bar(sl, title, subtitle=None):
    rect(sl, 0, 0, 13.333, 0.08, fill=TEAL)
    txt(sl, title, 0.45, 0.15, 12.5, 0.7, size=Pt(28), bold=True, color=NAVY)
    if subtitle:
        txt(sl, subtitle, 0.45, 0.82, 12.5, 0.38, size=Pt(13), italic=True, color=TEXT_MID)

def info_box(sl, title, body, l, t, w, h, title_color=TEAL, body_color=TEXT_MID, bg=WHITE, border=TEAL):
    rect(sl, l, t, w, h, fill=bg, line=border, lw=Pt(1))
    rect(sl, l, t, 0.05, h, fill=border)
    txt(sl, title, l + 0.12, t + 0.08, w - 0.18, 0.32, size=Pt(11), bold=True, color=title_color)
    txt(sl, body, l + 0.12, t + 0.40, w - 0.18, h - 0.5, size=Pt(10), color=body_color)

def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

prs = Presentation(r'C:\Users\pedro\OneDrive\Documentos\TCC_python\TCC_Apresentacao_Final.pptx')

# Adiciona o novo slide (vai pro final)
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg = sl.background.fill
bg.solid()
bg.fore_color.rgb = BG

header_bar(sl, "Base de Dados: Principais Variáveis", "Frequência mensal (2003–2026). Transformações em ln x 100 para leitura direta de elasticidades.")

# Coluna 1: O Choque e A Transmissão
rect(sl, 0.4, 1.4, 6.0, 5.5, fill=WHITE, line=DIVIDER, lw=Pt(1))
rect(sl, 0.4, 1.4, 6.0, 0.06, fill=TEAL)
txt(sl, "O Choque e A Cadeia de Transmissão", 0.55, 1.5, 5.7, 0.45, size=Pt(14), bold=True, color=TEAL)

info_box(sl, "Petróleo Brent (Variável de Choque)", "Preço global do barril de petróleo em USD. Fonte: IPEADATA/Banco Mundial.", 
         0.6, 2.0, 5.6, 0.9, bg=TEAL_LT, border=TEAL)

info_box(sl, "Preços de Combustíveis", "Gasolina A (refinaria), Gasolina C (bomba), Óleo Diesel e Etanol. Fonte: Agência Nacional do Petróleo (ANP).", 
         0.6, 3.1, 5.6, 0.9, bg=WHITE, border=TEAL)

info_box(sl, "Inflação (Variáveis Dependentes Finais)", "Índice Nacional de Preços ao Consumidor Amplo (IPCA Geral) e o subgrupo IPCA Transportes. Fonte: IBGE.", 
         0.6, 4.2, 5.6, 0.9, bg=WHITE, border=GREEN, title_color=GREEN)


# Coluna 2: Controles Macroeconômicos
rect(sl, 6.6, 1.4, 6.3, 5.5, fill=WHITE, line=DIVIDER, lw=Pt(1))
rect(sl, 6.6, 1.4, 6.3, 0.06, fill=NAVY)
txt(sl, "Controles Macroeconômicos (Ceteris Paribus)", 6.75, 1.5, 6.0, 0.45, size=Pt(14), bold=True, color=NAVY)

info_box(sl, "Taxa de Câmbio (R$/USD)", "Filtro primário que pode amplificar ou atenuar o choque externo em moeda local. Fonte: Banco Central.", 
         6.8, 2.0, 5.9, 0.8, bg=WHITE, border=NAVY, title_color=NAVY)

info_box(sl, "Atividade Econômica (IBC-Br)", "Proxy mensal do PIB para controlar choques de demanda doméstica por combustíveis. Fonte: Banco Central.", 
         6.8, 3.0, 5.9, 0.8, bg=WHITE, border=NAVY, title_color=NAVY)

info_box(sl, "Taxa Selic e Expectativas (Focus)", "Controle do ambiente de política monetária e antecipação inflacionária dos agentes econômicos.", 
         6.8, 4.0, 5.9, 0.8, bg=WHITE, border=NAVY, title_color=NAVY)

info_box(sl, "Índice Kilian (Demanda Global)", "Indicador de atividade econômica global. Isola o efeito de choques globais de demanda vs choques de oferta de petróleo.", 
         6.8, 5.0, 5.9, 0.8, bg=WHITE, border=GOLD, title_color=GOLD)

# Footer placeholder 
rect(sl, 0, 7.28, 13.333, 0.01, fill=DIVIDER)
txt(sl, "Pedro Franck Minella  |  TCC Ibmec-DF  |  2026", 0.45, 7.3, 9, 0.22, size=Pt(8), color=TEXT_LITE)

# O slide foi criado no final (índice len - 1). Quero que ele seja o slide 7 (índice 6), logo após "A Lente Metodológica"
total = len(prs.slides)
move_slide(prs, total - 1, 6)

# Atualizar todos os footers para ter "i/total" e limpar o antigo se der
for i, slide in enumerate(prs.slides):
    # Achar e remover o textbox que parece ser X/16 ou algo assim, e inserir novo
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            if shape.text.endswith(str(total-1)) or shape.text.endswith(str(total)):
                if "/" in shape.text and len(shape.text) < 8:
                    sp = shape._element
                    sp.getparent().remove(sp)
    txt(slide, f"{i+1}/{total}", 12.2, 7.3, 1.0, 0.22, size=Pt(8), color=TEXT_LITE, align=PP_ALIGN.RIGHT)

prs.save(r'C:\Users\pedro\OneDrive\Documentos\TCC_python\TCC_Apresentacao_Final.pptx')
print("Slide de variáveis adicionado com sucesso.")
