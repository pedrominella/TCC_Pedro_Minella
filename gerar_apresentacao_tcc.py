# -*- coding: utf-8 -*-
"""
gerar_apresentacao_tcc.py
Script para gerar a apresentação de defesa do TCC de Pedro Franck Minella.
Design: Dark academic com gradiente azul marinho institucional, fontes modernas
e layout premium para banca de graduação.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy
import os
from pathlib import Path

# ──────────────────────────────────────────
# PALETA DE CORES INSTITUCIONAL
# ──────────────────────────────────────────
NAVY       = RGBColor(0x10, 0x1C, 0x3E)   # Azul marinho fundo principal
DARK_BLUE  = RGBColor(0x1A, 0x2D, 0x5A)   # Azul escuro secundário
ACCENT     = RGBColor(0x2C, 0x82, 0xC9)   # Azul médio para destaques
GOLD       = RGBColor(0xF5, 0xC5, 0x18)   # Dourado para títulos especiais
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)   # Branco principal
LIGHT_GRAY = RGBColor(0xD6, 0xE0, 0xF0)   # Cinza claro para texto secundário
MID_GRAY   = RGBColor(0x8A, 0xA3, 0xC4)   # Cinza médio para bullets
GREEN_OK   = RGBColor(0x2E, 0xCC, 0x71)   # Verde para confirmações
RED_WARN   = RGBColor(0xE7, 0x4C, 0x3C)   # Vermelho para atenção
TRANS_BOX  = RGBColor(0x1E, 0x3A, 0x6E)   # Caixa de conteúdo

# ──────────────────────────────────────────
# CONFIGURAÇÃO DA APRESENTAÇÃO
# ──────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color: RGBColor):
    """Define a cor de fundo de um slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=Pt(0),
             transparency=0.0):
    """Adiciona um retângulo formatado ao slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()

    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()

    shape.shadow.inherit = False
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_name="Calibri", font_size=Pt(18),
                 bold=False, italic=False, color=WHITE,
                 align=PP_ALIGN.LEFT, word_wrap=True,
                 line_spacing=None):
    """Adiciona uma caixa de texto formatada ao slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if line_spacing:
        p.line_spacing = line_spacing
    return txBox


def add_paragraph_to_frame(tf, text, font_name="Calibri",
                            font_size=Pt(14), bold=False, italic=False,
                            color=WHITE, align=PP_ALIGN.LEFT,
                            space_before=Pt(6), bullet_char=None):
    """Adiciona um parágrafo formatado a um text_frame existente."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before

    if bullet_char:
        run0 = p.add_run()
        run0.text = bullet_char + "  "
        run0.font.name = font_name
        run0.font.size = font_size
        run0.font.bold = False
        run0.font.color.rgb = ACCENT

    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def add_accent_bar(slide, y_top=Inches(0.08), height=Inches(0.07)):
    """Adiciona barra de destaque dourada no topo do slide."""
    bar = add_rect(slide,
                   left=Inches(0), top=y_top,
                   width=SLIDE_W, height=height,
                   fill_color=GOLD)
    return bar


def add_footer(slide, slide_num, total=14):
    """Rodapé padronizado com número de slide e nome do autor."""
    # Linha fina
    add_rect(slide,
             left=Inches(0.5), top=Inches(7.0),
             width=Inches(12.33), height=Pt(1),
             fill_color=MID_GRAY)

    add_text_box(slide, "Pedro Franck Minella  |  TCC Ibmec-DF  |  2026",
                 left=Inches(0.5), top=Inches(7.05),
                 width=Inches(9), height=Inches(0.35),
                 font_size=Pt(9), color=MID_GRAY)

    add_text_box(slide, f"{slide_num}/{total}",
                 left=Inches(11.8), top=Inches(7.05),
                 width=Inches(1.0), height=Inches(0.35),
                 font_size=Pt(9), color=MID_GRAY,
                 align=PP_ALIGN.RIGHT)


def add_slide_title(slide, title_text, subtitle_text=None,
                    title_y=Inches(0.28), title_size=Pt(28)):
    """Adiciona título e subtítulo padronizados ao slide."""
    add_text_box(slide, title_text,
                 left=Inches(0.5), top=title_y,
                 width=Inches(12.0), height=Inches(0.7),
                 font_name="Calibri", font_size=title_size,
                 bold=True, color=WHITE)

    if subtitle_text:
        add_text_box(slide, subtitle_text,
                     left=Inches(0.5), top=title_y + Inches(0.6),
                     width=Inches(12.0), height=Inches(0.4),
                     font_name="Calibri", font_size=Pt(14),
                     italic=True, color=ACCENT)


def content_box(slide, left, top, width, height):
    """Cria uma caixa de conteúdo estilizada (fundo levemente mais claro)."""
    rect = add_rect(slide, left, top, width, height,
                    fill_color=TRANS_BOX)
    # Borda esquerda colorida
    add_rect(slide, left, top, Pt(4), height, fill_color=ACCENT)
    return rect


# ══════════════════════════════════════════════════════════════════════
# CONSTRUÇÃO DOS SLIDES
# ══════════════════════════════════════════════════════════════════════

def slide_01_capa(prs):
    """Slide 1: Capa Institucional"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide, y_top=Inches(0), height=Inches(0.12))

    # Logo/nome institution area
    add_rect(slide,
             left=Inches(0), top=Inches(0.12),
             width=Inches(13.333), height=Inches(1.2),
             fill_color=DARK_BLUE)

    add_text_box(slide, "INSTITUTO BRASILEIRO DE MERCADO DE CAPITAIS — IBMEC-DF",
                 left=Inches(0.5), top=Inches(0.25),
                 width=Inches(12.0), height=Inches(0.5),
                 font_name="Calibri", font_size=Pt(13),
                 bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Bacharelado em Ciências Econômicas",
                 left=Inches(0.5), top=Inches(0.72),
                 width=Inches(12.0), height=Inches(0.4),
                 font_name="Calibri", font_size=Pt(11),
                 color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Linha dourada decorativa
    add_rect(slide,
             left=Inches(2.5), top=Inches(1.5),
             width=Inches(8.33), height=Pt(2),
             fill_color=GOLD)

    # Título principal
    add_text_box(slide,
                 "Efeitos dos Choques do Preço do Petróleo\nsobre a Inflação no Brasil",
                 left=Inches(0.5), top=Inches(1.65),
                 width=Inches(12.33), height=Inches(1.4),
                 font_name="Calibri", font_size=Pt(32),
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "Evidências via Combustíveis, Petrobras e Local Projections (2003–2026)",
                 left=Inches(0.5), top=Inches(3.1),
                 width=Inches(12.33), height=Inches(0.5),
                 font_name="Calibri", font_size=Pt(15),
                 italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_rect(slide,
             left=Inches(2.5), top=Inches(3.75),
             width=Inches(8.33), height=Pt(1),
             fill_color=MID_GRAY)

    # Autor e orientador
    add_text_box(slide, "Pedro Franck Minella",
                 left=Inches(0.5), top=Inches(3.95),
                 width=Inches(12.33), height=Inches(0.5),
                 font_name="Calibri", font_size=Pt(18),
                 bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Orientador: Prof. Silvio Costa   |   Coordenador: Prof. Frederico Dias",
                 left=Inches(0.5), top=Inches(4.48),
                 width=Inches(12.33), height=Inches(0.4),
                 font_name="Calibri", font_size=Pt(12),
                 color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Ano
    add_text_box(slide, "Brasília, 2026",
                 left=Inches(0.5), top=Inches(5.0),
                 width=Inches(12.33), height=Inches(0.4),
                 font_name="Calibri", font_size=Pt(12),
                 color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Caixas de métricas na base
    metrics = [
        ("23 anos", "Amostra Temporal\n(2003–2026)"),
        ("4 modelos", "Especificações\nEconométricas"),
        ("6 variáveis", "Preços de\nCombustíveis"),
        ("12 meses", "Horizonte de\nProjeção"),
    ]
    box_w = Inches(2.8)
    box_h = Inches(1.1)
    for i, (val, lbl) in enumerate(metrics):
        bx = Inches(0.4) + i * Inches(3.15)
        add_rect(slide, left=bx, top=Inches(6.05),
                 width=box_w, height=box_h,
                 fill_color=DARK_BLUE)
        add_rect(slide, left=bx, top=Inches(6.05),
                 width=Pt(3), height=box_h,
                 fill_color=ACCENT)
        add_text_box(slide, val,
                     left=bx + Inches(0.15), top=Inches(6.08),
                     width=box_w - Inches(0.2), height=Inches(0.45),
                     font_size=Pt(18), bold=True, color=GOLD)
        add_text_box(slide, lbl,
                     left=bx + Inches(0.15), top=Inches(6.5),
                     width=box_w - Inches(0.2), height=Inches(0.55),
                     font_size=Pt(9), color=LIGHT_GRAY)


def slide_02_mecanismo(prs):
    """Slide 2: Mecanismo de Transmissão"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 2)

    add_slide_title(slide, "O Mecanismo de Transmissão",
                    "Como o choque global chega ao consumidor brasileiro?")

    # Fluxograma horizontal
    nodes = [
        ("🛢️", "BRENT\n(USD)", ACCENT),
        ("💱", "CÂMBIO\n(R$/USD)", ACCENT),
        ("🏭", "PETROBRAS\n(Refinaria)", GOLD),
        ("⛽", "POSTO\n(Bomba)", ACCENT),
        ("📊", "IPCA\nTransportes", GREEN_OK),
    ]

    node_w = Inches(2.0)
    node_h = Inches(1.5)
    y_node = Inches(1.7)
    spacing = Inches(2.2)
    start_x = Inches(0.45)

    for i, (icon, label, color) in enumerate(nodes):
        nx = start_x + i * spacing
        # Caixa do nó
        add_rect(slide, left=nx, top=y_node,
                 width=node_w, height=node_h,
                 fill_color=DARK_BLUE)
        add_rect(slide, left=nx, top=y_node,
                 width=node_w, height=Pt(4),
                 fill_color=color)
        # Ícone
        add_text_box(slide, icon,
                     left=nx, top=y_node + Inches(0.05),
                     width=node_w, height=Inches(0.55),
                     font_size=Pt(22), align=PP_ALIGN.CENTER)
        # Label
        add_text_box(slide, label,
                     left=nx, top=y_node + Inches(0.6),
                     width=node_w, height=Inches(0.85),
                     font_size=Pt(11), bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        # Seta
        if i < len(nodes) - 1:
            add_text_box(slide, "→",
                         left=nx + node_w + Inches(0.05),
                         top=y_node + Inches(0.5),
                         width=Inches(0.3), height=Inches(0.5),
                         font_size=Pt(22), bold=True, color=ACCENT,
                         align=PP_ALIGN.CENTER)

    # Caixa IPCA Geral (secundária, abaixo)
    add_rect(slide,
             left=Inches(10.8), top=Inches(3.55),
             width=Inches(2.1), height=Inches(1.1),
             fill_color=DARK_BLUE)
    add_rect(slide,
             left=Inches(10.8), top=Inches(3.55),
             width=Inches(2.1), height=Pt(3),
             fill_color=MID_GRAY)
    add_text_box(slide, "📊 IPCA\nGeral",
                 left=Inches(10.8), top=Inches(3.6),
                 width=Inches(2.1), height=Inches(0.95),
                 font_size=Pt(11), color=LIGHT_GRAY,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, "↑ Efeito concentrado\ne persistente (43%)",
                 left=Inches(10.5), top=Inches(3.3),
                 width=Inches(2.5), height=Inches(0.4),
                 font_size=Pt(9), color=GREEN_OK,
                 align=PP_ALIGN.CENTER)

    add_text_box(slide, "↑ Efeito diluído\ne transitório",
                 left=Inches(10.5), top=Inches(4.7),
                 width=Inches(2.5), height=Inches(0.4),
                 font_size=Pt(9), color=MID_GRAY,
                 align=PP_ALIGN.CENTER)

    # Notas contextuais abaixo dos nós
    notas = [
        "Preço int.\nem dólares",
        "Amplifica\nou atenua\no choque",
        "Define timing\ne magnitude\ndo repasse",
        "Impostos,\nmistura de\netanol",
        "Peso direto\nna cesta do\nconsumidor",
    ]
    for i, nota in enumerate(notas):
        nx = start_x + i * spacing
        add_text_box(slide, nota,
                     left=nx, top=Inches(3.35),
                     width=node_w, height=Inches(0.7),
                     font_size=Pt(9), italic=True, color=MID_GRAY,
                     align=PP_ALIGN.CENTER)

    # Box de destaque na base
    add_rect(slide,
             left=Inches(0.4), top=Inches(4.25),
             width=Inches(12.5), height=Inches(1.1),
             fill_color=DARK_BLUE)
    add_rect(slide,
             left=Inches(0.4), top=Inches(4.25),
             width=Pt(4), height=Inches(1.1),
             fill_color=GOLD)
    add_text_box(slide,
                 "Hipótese Central: O repasse é setorialmente concentrado no IPCA Transportes e a política de preços "
                 "da Petrobras atua como válvula reguladora da velocidade e da magnitude dessa transmissão.",
                 left=Inches(0.65), top=Inches(4.32),
                 width=Inches(12.0), height=Inches(0.95),
                 font_size=Pt(11), italic=True, color=LIGHT_GRAY)

    # Caixas de papel institucional
    roles = [
        ("Filtro Cambial", "Converte o choque em\nReais, amplificando\nvolatilidade doméstica"),
        ("Filtro Petrobras", "Decisão corporativa de\nrepasse; o pilar central\ndesta pesquisa"),
        ("Filtro de Mercado", "Distribuição, impostos e\nsubstituição etanol/gasolina"),
    ]
    box_w2 = Inches(3.9)
    for i, (titulo, desc) in enumerate(roles):
        bx2 = Inches(0.4) + i * Inches(4.3)
        add_rect(slide, left=bx2, top=Inches(5.5),
                 width=box_w2, height=Inches(1.4),
                 fill_color=DARK_BLUE)
        add_text_box(slide, titulo,
                     left=bx2 + Inches(0.1), top=Inches(5.52),
                     width=box_w2 - Inches(0.15), height=Inches(0.4),
                     font_size=Pt(11), bold=True, color=ACCENT)
        add_text_box(slide, desc,
                     left=bx2 + Inches(0.1), top=Inches(5.9),
                     width=box_w2 - Inches(0.15), height=Inches(0.9),
                     font_size=Pt(9.5), color=LIGHT_GRAY)


def slide_03_petrobras(prs):
    """Slide 3: A Mudança Institucional da Petrobras"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 3)

    add_slide_title(slide, "A Petrobras como Reguladora do Repasse",
                    "Setembro de 2016: A quebra estrutural da política de preços")

    # Duas colunas: Pré e Pós 2016
    col_w = Inches(5.8)
    col_h = Inches(4.5)
    y_col = Inches(1.55)

    # Pré 2016
    add_rect(slide, left=Inches(0.3), top=y_col,
             width=col_w, height=col_h, fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(0.3), top=y_col,
             width=col_w, height=Pt(5), fill_color=MID_GRAY)
    add_text_box(slide, "🛡️  PRÉ-SETEMBRO/2016",
                 left=Inches(0.4), top=y_col + Inches(0.1),
                 width=col_w - Inches(0.15), height=Inches(0.45),
                 font_size=Pt(14), bold=True, color=LIGHT_GRAY)
    add_text_box(slide, "Controle Discricionário de Preços",
                 left=Inches(0.4), top=y_col + Inches(0.55),
                 width=col_w - Inches(0.15), height=Inches(0.35),
                 font_size=Pt(11), italic=True, color=MID_GRAY)

    pre_bullets = [
        "Preços internos descolados do Brent internacional",
        "Governo usava a estatal como instrumento anti-inflacionário",
        "Defasagens de meses entre choque externo e reajuste interno",
        "Resultado: repasse inflacionário BLOQUEADO",
        "Custo: acúmulo de passivo corporativo bilionário na Petrobras",
    ]
    for i, b in enumerate(pre_bullets):
        color = RED_WARN if "BLOQUEADO" in b else LIGHT_GRAY
        bold = "BLOQUEADO" in b
        add_text_box(slide, f"  ›  {b}",
                     left=Inches(0.4), top=y_col + Inches(1.05) + i * Inches(0.62),
                     width=col_w - Inches(0.15), height=Inches(0.58),
                     font_size=Pt(10.5), bold=bold, color=color)

    # Pós 2016
    add_rect(slide, left=Inches(7.2), top=y_col,
             width=col_w, height=col_h, fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(7.2), top=y_col,
             width=col_w, height=Pt(5), fill_color=GOLD)
    add_text_box(slide, "⚡  PÓS-SETEMBRO/2016 (PPI)",
                 left=Inches(7.3), top=y_col + Inches(0.1),
                 width=col_w - Inches(0.15), height=Inches(0.45),
                 font_size=Pt(14), bold=True, color=GOLD)
    add_text_box(slide, "Política de Paridade de Importação (PPI)",
                 left=Inches(7.3), top=y_col + Inches(0.55),
                 width=col_w - Inches(0.15), height=Inches(0.35),
                 font_size=Pt(11), italic=True, color=MID_GRAY)

    pos_bullets = [
        "Preços domésticos indexados ao Brent e ao câmbio",
        "Reajustes frequentes acompanhando o mercado externo",
        "Petrobras deixa de ser o colchão de estabilidade inflacionária",
        "Resultado: repasse inflacionário ATIVADO e IMEDIATO",
        "Brasil passa a ser mais exposto a tensões geopolíticas globais",
    ]
    for i, b in enumerate(pos_bullets):
        color = GREEN_OK if "ATIVADO" in b else LIGHT_GRAY
        bold = "ATIVADO" in b
        add_text_box(slide, f"  ›  {b}",
                     left=Inches(7.3), top=y_col + Inches(1.05) + i * Inches(0.62),
                     width=col_w - Inches(0.15), height=Inches(0.58),
                     font_size=Pt(10.5), bold=bold, color=color)

    # Seta central
    add_text_box(slide, "→",
                 left=Inches(6.1), top=y_col + Inches(1.8),
                 width=Inches(1.0), height=Inches(0.7),
                 font_size=Pt(36), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide, "2016",
                 left=Inches(6.0), top=y_col + Inches(2.5),
                 width=Inches(1.2), height=Inches(0.45),
                 font_size=Pt(13), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Caixa de hipótese
    add_rect(slide,
             left=Inches(0.3), top=Inches(6.2),
             width=Inches(12.7), height=Inches(1.0),
             fill_color=DARK_BLUE)
    add_rect(slide,
             left=Inches(0.3), top=Inches(6.2),
             width=Pt(4), height=Inches(1.0),
             fill_color=GOLD)
    add_text_box(slide,
                 "Hipótese Institucional Testada: A adoção do PPI aumentou significativamente "
                 "a velocidade e magnitude da transmissão do choque de petróleo para o IPCA Transportes.",
                 left=Inches(0.55), top=Inches(6.27),
                 width=Inches(12.15), height=Inches(0.85),
                 font_size=Pt(11), italic=True, color=LIGHT_GRAY)


def slide_04_perguntas(prs):
    """Slide 4: Perguntas de Pesquisa"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 4)

    add_slide_title(slide, "Perguntas de Pesquisa",
                    "Três questões econômicas fundamentais que este TCC responde")

    perguntas = [
        {
            "num": "01",
            "titulo": "Canal Upstream — Petróleo → Combustíveis",
            "body": ("Qual é a magnitude e a velocidade do repasse do Brent ao preço da "
                     "gasolina, do diesel e do etanol no Brasil? O impacto é rápido e permanente "
                     "ou gradual e temporário?"),
            "cor": ACCENT,
        },
        {
            "num": "02",
            "titulo": "Canal Downstream — Combustíveis → Inflação",
            "body": ("O choque nos combustíveis contamina a inflação de forma ampla (IPCA Geral) "
                     "ou o impacto é setorialmente concentrado no componente de Transportes?"),
            "cor": GOLD,
        },
        {
            "num": "03",
            "titulo": "Efeito Institucional — O Papel da Petrobras",
            "body": ("A mudança para o regime PPI em setembro de 2016 alterou estruturalmente "
                     "a velocidade e magnitude do repasse inflacionário dos combustíveis para o "
                     "IPCA Transportes?"),
            "cor": GREEN_OK,
        },
    ]

    for i, pq in enumerate(perguntas):
        box_y = Inches(1.6) + i * Inches(1.7)
        add_rect(slide, left=Inches(0.4), top=box_y,
                 width=Inches(12.5), height=Inches(1.55),
                 fill_color=DARK_BLUE)
        add_rect(slide, left=Inches(0.4), top=box_y,
                 width=Inches(0.7), height=Inches(1.55),
                 fill_color=pq["cor"])
        add_text_box(slide, pq["num"],
                     left=Inches(0.4), top=box_y + Inches(0.4),
                     width=Inches(0.7), height=Inches(0.55),
                     font_size=Pt(22), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(slide, pq["titulo"],
                     left=Inches(1.25), top=box_y + Inches(0.1),
                     width=Inches(11.3), height=Inches(0.45),
                     font_size=Pt(13), bold=True, color=pq["cor"])
        add_text_box(slide, pq["body"],
                     left=Inches(1.25), top=box_y + Inches(0.55),
                     width=Inches(11.3), height=Inches(0.9),
                     font_size=Pt(11), color=LIGHT_GRAY)


def slide_05_metodologia(prs):
    """Slide 5: Por que Projeções Locais?"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 5)

    add_slide_title(slide, "Metodologia: Projeções Locais (Jordà, 2005)",
                    "Por que Local Projections em vez do VAR clássico?")

    # Dois painéis
    # Painel esquerdo: LP
    add_rect(slide, left=Inches(0.3), top=Inches(1.55),
             width=Inches(6.1), height=Inches(4.0),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(0.3), top=Inches(1.55),
             width=Inches(6.1), height=Pt(5), fill_color=ACCENT)
    add_text_box(slide, "✅  LOCAL PROJECTIONS (LP)",
                 left=Inches(0.4), top=Inches(1.6),
                 width=Inches(5.9), height=Inches(0.45),
                 font_size=Pt(13), bold=True, color=ACCENT)

    lp_points = [
        "Estima regressão OLS direta para cada horizonte h",
        "Não depende de restrições dinâmicas do sistema",
        "Facilita introdução de regimes e não-linearidades",
        "Robusto a erros de especificação do VAR",
        "Erros-padrão HAC (Newey-West) corrigem autocorrelação por construção",
        "Padrão moderno na macroeconomia empírica (Ramey, 2016; Stock & Watson, 2018)",
    ]
    for i, pt in enumerate(lp_points):
        add_text_box(slide, f"  ✓  {pt}",
                     left=Inches(0.4), top=Inches(2.15) + i * Inches(0.52),
                     width=Inches(5.9), height=Inches(0.48),
                     font_size=Pt(10.5), color=LIGHT_GRAY)

    # Painel direito: VAR
    add_rect(slide, left=Inches(6.9), top=Inches(1.55),
             width=Inches(6.1), height=Inches(4.0),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(6.9), top=Inches(1.55),
             width=Inches(6.1), height=Pt(5), fill_color=MID_GRAY)
    add_text_box(slide, "⚠️  VAR (para referência e robustez)",
                 left=Inches(7.0), top=Inches(1.6),
                 width=Inches(5.9), height=Inches(0.45),
                 font_size=Pt(13), bold=True, color=MID_GRAY)

    var_points = [
        "Itera respostas através de coeficientes de curto prazo",
        "Amplifica erros de especificação ao longo do tempo",
        "Restrições de identificação podem ser arbitrárias",
        "Menos flexível para modelar quebras estruturais",
        "Menos adequado para análise de regimes não-lineares",
        "Utilizado como robustez no Apêndice A deste trabalho",
    ]
    for i, pt in enumerate(var_points):
        color = LIGHT_GRAY if i < 5 else ACCENT
        add_text_box(slide, f"  –  {pt}",
                     left=Inches(7.0), top=Inches(2.15) + i * Inches(0.52),
                     width=Inches(5.9), height=Inches(0.48),
                     font_size=Pt(10.5), color=color)

    # Equação simplificada
    add_rect(slide,
             left=Inches(0.3), top=Inches(5.7),
             width=Inches(12.7), height=Inches(1.5),
             fill_color=DARK_BLUE)
    add_rect(slide,
             left=Inches(0.3), top=Inches(5.7),
             width=Pt(4), height=Inches(1.5),
             fill_color=ACCENT)
    add_text_box(slide, "Especificação Central:  y(t+h) = α(h) + β(h) · ΔPetróleo(t) + γ(h)·Controles(t) + ε(t+h)",
                 left=Inches(0.55), top=Inches(5.78),
                 width=Inches(12.2), height=Inches(0.5),
                 font_size=Pt(12), bold=True, color=WHITE)
    add_text_box(slide,
                 "Para cada h ∈ {0,1,...,12}: estimação independente por OLS com erros HAC(maxlags=h). "
                 "Controles: câmbio, IBC-Br, Selic, expectativas Focus e Índice Kilian de demanda global.",
                 left=Inches(0.55), top=Inches(6.28),
                 width=Inches(12.2), height=Inches(0.8),
                 font_size=Pt(10), italic=True, color=MID_GRAY)


def slide_06_estrategia(prs):
    """Slide 6: Estratégia Empírica e Controles"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 6)

    add_slide_title(slide, "Estratégia Empírica e Controles",
                    "Conjunto de variáveis e blindagens contra correlações espúrias")

    # Coluna esquerda: variáveis do modelo
    add_rect(slide, left=Inches(0.3), top=Inches(1.55),
             width=Inches(5.9), height=Inches(5.6),
             fill_color=DARK_BLUE)
    add_text_box(slide, "Variáveis do Modelo",
                 left=Inches(0.4), top=Inches(1.6),
                 width=Inches(5.7), height=Inches(0.45),
                 font_size=Pt(13), bold=True, color=GOLD)

    vars_model = [
        ("CHOQUE DE INTERESSE", ACCENT, [
            "Variação % mensal do Brent (USD e BRL)",
            "Transformação: Δln × 100",
        ]),
        ("VARIÁVEIS RESPOSTA", ACCENT, [
            "Gasolina C (bomba), Gasolina A (refinaria)",
            "Óleo Diesel, Etanol",
            "IPCA Geral e IPCA Transportes",
        ]),
        ("CONTROLES DINÂMICOS", MID_GRAY, [
            "Taxa de câmbio R$/USD (contemporâneo + lags)",
            "IBC-Br — proxy de atividade econômica",
            "Taxa Selic (política monetária)",
            "Expectativa de inflação Focus",
            "Índice Kilian — demanda global de commodities",
            "Dummies sazonais mensais (11 dummies)",
        ]),
    ]

    y_off = Inches(2.15)
    for grupo, cor, items in vars_model:
        add_text_box(slide, grupo,
                     left=Inches(0.45), top=y_off,
                     width=Inches(5.65), height=Inches(0.35),
                     font_size=Pt(9.5), bold=True, color=cor)
        y_off += Inches(0.35)
        for it in items:
            add_text_box(slide, f"  ·  {it}",
                         left=Inches(0.45), top=y_off,
                         width=Inches(5.65), height=Inches(0.38),
                         font_size=Pt(9.5), color=LIGHT_GRAY)
            y_off += Inches(0.38)
        y_off += Inches(0.12)

    # Coluna direita: lógica de identificação
    add_rect(slide, left=Inches(6.6), top=Inches(1.55),
             width=Inches(6.4), height=Inches(5.6),
             fill_color=DARK_BLUE)
    add_text_box(slide, "Lógica de Identificação",
                 left=Inches(6.7), top=Inches(1.6),
                 width=Inches(6.2), height=Inches(0.45),
                 font_size=Pt(13), bold=True, color=GOLD)

    logica = [
        ("Por que controlar o câmbio?",
         "O R$/USD reage a choques fiscais e políticos domésticos. Incluir o câmbio "
         "como controle isola o efeito do preço do barril de ruídos macroeconômicos internos."),
        ("Por que o Índice Kilian?",
         "O Brent sobe tanto quando há booms de demanda global como quando há choques de "
         "oferta. O Índice Kilian capta a demanda por fretes secos e, ao controlá-lo, "
         "isolamos o efeito do petróleo independente do ciclo global."),
        ("Por que 3 defasagens?",
         "Defasagens da variável dependente, do choque e dos controles absorvem dinâmicas "
         "autorregressivas e filtram memória de curto prazo nas séries mensais."),
        ("Correção HAC (Newey-West)",
         "Local Projections acumulam autocorrelação MA(h) por construção. "
         "A correção Newey-West com maxlags=h garante erros-padrão corretos e bandas "
         "de confiança precisas, sem inflação artificial de significância."),
    ]

    y2 = Inches(2.15)
    for titulo, desc in logica:
        add_text_box(slide, titulo,
                     left=Inches(6.75), top=y2,
                     width=Inches(6.1), height=Inches(0.35),
                     font_size=Pt(10.5), bold=True, color=ACCENT)
        add_text_box(slide, desc,
                     left=Inches(6.75), top=y2 + Inches(0.35),
                     width=Inches(6.1), height=Inches(0.75),
                     font_size=Pt(9.5), color=LIGHT_GRAY)
        y2 += Inches(1.2)


def slide_07_combustiveis(prs):
    """Slide 7: Resultados — Petróleo → Combustíveis"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 7)

    add_slide_title(slide, "Resultados: Petróleo → Combustíveis",
                    "O primeiro elo da cadeia está estatisticamente ativo")

    # Tabela de resultados
    headers = ["Horizonte", "Diesel", "Gasolina Refinaria", "Gasolina Bomba", "Etanol"]
    rows_data = [
        ["h = 0", "0,115 ***", "0,154 ***", "0,269 ***", "0,098 ***"],
        ["h = 1", "0,166 ***", "0,295 ***", "0,397 ***", "0,139 ***"],
        ["h = 2", "0,150 *",   "0,321 ***", "0,437 ***", "0,137 ***"],
        ["h = 3", "0,105",     "0,353 ***", "0,434 ***", "0,129 ***"],
        ["h = 6", "−0,044",    "0,323 ***", "0,464 ***", "0,178 ***"],
    ]

    col_w_table = [Inches(1.1), Inches(1.1), Inches(2.0), Inches(2.0), Inches(1.1)]
    row_h = Inches(0.45)
    table_left = Inches(0.4)
    table_top = Inches(1.65)

    # Cabeçalho
    x = table_left
    for j, h in enumerate(headers):
        add_rect(slide, left=x, top=table_top,
                 width=col_w_table[j], height=row_h,
                 fill_color=ACCENT)
        add_text_box(slide, h,
                     left=x + Inches(0.05), top=table_top + Inches(0.06),
                     width=col_w_table[j] - Inches(0.07), height=row_h - Inches(0.1),
                     font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += col_w_table[j]

    for i, row in enumerate(rows_data):
        y_row = table_top + (i + 1) * row_h
        x = table_left
        row_fill = DARK_BLUE if i % 2 == 0 else RGBColor(0x14, 0x25, 0x4D)
        for j, cell in enumerate(row):
            add_rect(slide, left=x, top=y_row,
                     width=col_w_table[j], height=row_h, fill_color=row_fill)
            # Color based on significance
            if "***" in cell:
                tc = GREEN_OK
            elif "*" in cell:
                tc = GOLD
            elif "−" in cell or cell.strip() == "—":
                tc = RED_WARN
            elif j == 0:
                tc = LIGHT_GRAY
            else:
                tc = MID_GRAY
            bold_cell = j == 0 or "***" in cell
            add_text_box(slide, cell,
                         left=x + Inches(0.05), top=y_row + Inches(0.06),
                         width=col_w_table[j] - Inches(0.07), height=row_h - Inches(0.1),
                         font_size=Pt(10), bold=bold_cell, color=tc,
                         align=PP_ALIGN.CENTER)
            x += col_w_table[j]

    # Legenda
    add_text_box(slide, "*** p<0,01   * p<0,10   Coeficientes acumulados em p.p. por 1 p.p. de choque no Brent",
                 left=Inches(0.4), top=table_top + 6 * row_h + Inches(0.1),
                 width=Inches(7.8), height=Inches(0.35),
                 font_size=Pt(9), italic=True, color=MID_GRAY)

    # Insights à direita
    add_rect(slide, left=Inches(8.0), top=Inches(1.65),
             width=Inches(5.0), height=Inches(4.2),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(8.0), top=Inches(1.65),
             width=Pt(4), height=Inches(4.2), fill_color=GOLD)
    add_text_box(slide, "Insights Econômicos",
                 left=Inches(8.1), top=Inches(1.7),
                 width=Inches(4.8), height=Inches(0.4),
                 font_size=Pt(12), bold=True, color=GOLD)

    insights = [
        ("⚡ Diesel", "Repasse ultra-rápido (h=0, h=1) e temporário — canal logístico puro de custo de frete"),
        ("⛽ Gasolina Bomba", "Significância mantida em todos os horizontes — canal direto ao consumidor final"),
        ("🏭 Gasolina Refinaria", "Acumulação gradual até h=4 — reflete janela decisória da Petrobras"),
        ("🌿 Etanol", "Resposta persistente (h>6) — lei de substituição nos veículos flex"),
    ]
    for k, (titulo, desc) in enumerate(insights):
        y_ins = Inches(2.2) + k * Inches(0.95)
        add_text_box(slide, titulo,
                     left=Inches(8.1), top=y_ins,
                     width=Inches(4.8), height=Inches(0.35),
                     font_size=Pt(10.5), bold=True, color=ACCENT)
        add_text_box(slide, desc,
                     left=Inches(8.1), top=y_ins + Inches(0.35),
                     width=Inches(4.8), height=Inches(0.5),
                     font_size=Pt(9.5), color=LIGHT_GRAY)

    # Box de destaque
    add_rect(slide, left=Inches(0.4), top=Inches(6.25),
             width=Inches(12.5), height=Inches(0.9),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(0.4), top=Inches(6.25),
             width=Pt(4), height=Inches(0.9), fill_color=GREEN_OK)
    add_text_box(slide,
                 "Confirmação do Canal Upstream: Todos os combustíveis respondem de forma positiva e estatisticamente "
                 "significante ao choque do petróleo. O primeiro elo da cadeia está comprovado.",
                 left=Inches(0.65), top=Inches(6.3),
                 width=Inches(12.1), height=Inches(0.8),
                 font_size=Pt(10.5), bold=True, color=WHITE)


def slide_08_gasolina_ipca(prs):
    """Slide 8: Gasolina C → IPCA Transportes vs. Geral"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 8)

    add_slide_title(slide, "Gasolina → IPCA Transportes vs. IPCA Geral",
                    "O repasse inflacionário é setorialmente concentrado")

    # Painel esquerdo: IPCA Transportes
    add_rect(slide, left=Inches(0.3), top=Inches(1.55),
             width=Inches(6.1), height=Inches(4.7),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(0.3), top=Inches(1.55),
             width=Inches(6.1), height=Pt(5), fill_color=GREEN_OK)
    add_text_box(slide, "IPCA TRANSPORTES  ✅",
                 left=Inches(0.4), top=Inches(1.6),
                 width=Inches(5.9), height=Inches(0.45),
                 font_size=Pt(13), bold=True, color=GREEN_OK)

    dados_transp = [
        ("h = 0", "0,269 ***", "repasse imediato de 26,9%"),
        ("h = 1", "0,397 ***", "aceleração no primeiro mês"),
        ("h = 2", "0,437 ***", "consolidação do repasse"),
        ("h = 3", "0,434 ***", "estabilização robusta"),
        ("h = 6", "0,464 ***", "mantém força no médio prazo"),
        ("h = 12", "0,431 ***", "persistente no longo prazo"),
    ]

    for i, (h, coef, desc) in enumerate(dados_transp):
        y_r = Inches(2.15) + i * Inches(0.52)
        add_text_box(slide, h,
                     left=Inches(0.4), top=y_r,
                     width=Inches(0.85), height=Inches(0.45),
                     font_size=Pt(10), bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, coef,
                     left=Inches(1.3), top=y_r,
                     width=Inches(1.3), height=Inches(0.45),
                     font_size=Pt(10), bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc,
                     left=Inches(2.7), top=y_r,
                     width=Inches(3.5), height=Inches(0.45),
                     font_size=Pt(9.5), color=MID_GRAY)

    # Barra de 43%
    add_rect(slide, left=Inches(0.4), top=Inches(5.35),
             width=Inches(6.1), height=Inches(0.7),
             fill_color=RGBColor(0x0D, 0x3D, 0x20))
    add_text_box(slide,
                 "Pass-through de longo prazo: ~43%  |  t-stat > 3,6 em todo o horizonte",
                 left=Inches(0.5), top=Inches(5.42),
                 width=Inches(5.9), height=Inches(0.55),
                 font_size=Pt(11), bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)

    # Painel direito: IPCA Geral
    add_rect(slide, left=Inches(6.9), top=Inches(1.55),
             width=Inches(6.1), height=Inches(4.7),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(6.9), top=Inches(1.55),
             width=Inches(6.1), height=Pt(5), fill_color=MID_GRAY)
    add_text_box(slide, "IPCA GERAL  ⚠️",
                 left=Inches(7.0), top=Inches(1.6),
                 width=Inches(5.9), height=Inches(0.45),
                 font_size=Pt(13), bold=True, color=MID_GRAY)

    dados_geral = [
        ("h = 0", "0,035 ***", "efeito inicial pequeno mas sig."),
        ("h = 1", "0,057 ***", "pico máximo logo dissipado"),
        ("h = 2", "0,055 ***", "mantém apenas a 5%"),
        ("h = 3", "0,047 **",  "significância reduzida"),
        ("h = 6", "0,040",     "sem significância estatística"),
        ("h = 12", "−0,013",   "sem significância, cruza zero"),
    ]

    for i, (h, coef, desc) in enumerate(dados_geral):
        y_r = Inches(2.15) + i * Inches(0.52)
        sig = "***" in coef or "**" in coef
        cor_c = GOLD if sig else RED_WARN
        add_text_box(slide, h,
                     left=Inches(7.0), top=y_r,
                     width=Inches(0.85), height=Inches(0.45),
                     font_size=Pt(10), bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, coef,
                     left=Inches(7.9), top=y_r,
                     width=Inches(1.3), height=Inches(0.45),
                     font_size=Pt(10), bold=True, color=cor_c, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc,
                     left=Inches(9.3), top=y_r,
                     width=Inches(3.5), height=Inches(0.45),
                     font_size=Pt(9.5), color=MID_GRAY)

    add_rect(slide, left=Inches(6.9), top=Inches(5.35),
             width=Inches(6.1), height=Inches(0.7),
             fill_color=RGBColor(0x2A, 0x1A, 0x0A))
    add_text_box(slide,
                 "Pass-through diluído:  < 0,06 p.p.  |  Insignificante a partir do 6.º mês",
                 left=Inches(7.0), top=Inches(5.42),
                 width=Inches(5.9), height=Inches(0.55),
                 font_size=Pt(11), bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Box de conclusão
    add_rect(slide, left=Inches(0.3), top=Inches(6.2),
             width=Inches(12.7), height=Inches(1.0),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(0.3), top=Inches(6.2),
             width=Pt(4), height=Inches(1.0), fill_color=GREEN_OK)
    add_text_box(slide,
                 "Conclusão: O impacto inflacionário da gasolina é 7× maior no IPCA Transportes que no IPCA Geral. "
                 "O efeito sobre a inflação geral é modesto e transitório — típico de choque de custo setorial.",
                 left=Inches(0.55), top=Inches(6.27),
                 width=Inches(12.15), height=Inches(0.85),
                 font_size=Pt(11), italic=True, color=WHITE)


def slide_09_diesel(prs):
    """Slide 9: Diesel — Impacto Logístico Temporário"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 9)

    add_slide_title(slide, "Óleo Diesel: Choque Logístico Temporário",
                    "Por que o impacto do diesel se dissipa em dois meses?")

    # Gráfico de barras estilizado manual (resposta acumulada)
    horizontes = [0, 1, 2, 3, 4, 5, 6]
    coefs_transp = [0.115, 0.166, 0.150, 0.105, 0.023, -0.015, -0.044]
    coefs_geral  = [0.010, 0.028, 0.016, 0.012, -0.003, 0.000, -0.019]

    chart_left = Inches(0.4)
    chart_top  = Inches(1.65)
    chart_w    = Inches(7.5)
    chart_h    = Inches(4.8)

    add_rect(slide, left=chart_left, top=chart_top,
             width=chart_w, height=chart_h, fill_color=DARK_BLUE)

    max_val = 0.22
    bar_w = Inches(0.45)
    group_w = Inches(0.95)
    y_baseline = chart_top + chart_h * 0.55  # linha zero

    for i, h in enumerate(horizontes):
        x_group = chart_left + Inches(0.5) + i * group_w

        # Barra Transportes
        v_transp = coefs_transp[i]
        bar_h_transp = abs(v_transp) / max_val * Inches(2.0)
        bar_y_transp = y_baseline - bar_h_transp if v_transp >= 0 else y_baseline
        col_t = GREEN_OK if v_transp > 0 else RED_WARN
        add_rect(slide, left=x_group, top=bar_y_transp,
                 width=bar_w * 0.6, height=bar_h_transp, fill_color=col_t)

        # Barra Geral
        v_geral = coefs_geral[i]
        bar_h_geral = abs(v_geral) / max_val * Inches(2.0)
        bar_y_geral = y_baseline - bar_h_geral if v_geral >= 0 else y_baseline
        col_g = LIGHT_GRAY if v_geral > 0 else RED_WARN
        add_rect(slide, left=x_group + bar_w * 0.65, top=bar_y_geral,
                 width=bar_w * 0.6, height=bar_h_geral, fill_color=col_g)

        # Label h
        add_text_box(slide, f"h={h}",
                     left=x_group, top=y_baseline + Inches(0.08),
                     width=bar_w * 1.3, height=Inches(0.3),
                     font_size=Pt(9), color=MID_GRAY, align=PP_ALIGN.CENTER)

        # Valor transportes
        add_text_box(slide, f"{v_transp:.3f}",
                     left=x_group, top=bar_y_transp - Inches(0.3),
                     width=bar_w * 0.8, height=Inches(0.28),
                     font_size=Pt(8), bold=True,
                     color=GREEN_OK if v_transp >= 0 else RED_WARN,
                     align=PP_ALIGN.CENTER)

    # Linha zero
    add_rect(slide, left=chart_left + Inches(0.3), top=y_baseline,
             width=chart_w - Inches(0.4), height=Pt(1),
             fill_color=WHITE)

    # Legenda do gráfico
    add_rect(slide, left=chart_left + Inches(0.3), top=chart_top + chart_h - Inches(0.9),
             width=Inches(0.25), height=Inches(0.25), fill_color=GREEN_OK)
    add_text_box(slide, "Diesel → IPCA Transportes",
                 left=chart_left + Inches(0.6), top=chart_top + chart_h - Inches(0.9),
                 width=Inches(2.5), height=Inches(0.3),
                 font_size=Pt(9), color=GREEN_OK)
    add_rect(slide, left=chart_left + Inches(3.2), top=chart_top + chart_h - Inches(0.9),
             width=Inches(0.25), height=Inches(0.25), fill_color=LIGHT_GRAY)
    add_text_box(slide, "Diesel → IPCA Geral",
                 left=chart_left + Inches(3.5), top=chart_top + chart_h - Inches(0.9),
                 width=Inches(2.5), height=Inches(0.3),
                 font_size=Pt(9), color=LIGHT_GRAY)

    # Painel direito de análise
    add_rect(slide, left=Inches(8.3), top=Inches(1.65),
             width=Inches(4.7), height=Inches(4.8),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(8.3), top=Inches(1.65),
             width=Pt(4), height=Inches(4.8), fill_color=ACCENT)
    add_text_box(slide, "Por que o Diesel se dissipa?",
                 left=Inches(8.4), top=Inches(1.7),
                 width=Inches(4.55), height=Inches(0.4),
                 font_size=Pt(12), bold=True, color=ACCENT)

    razoes = [
        ("Tarifas Administradas", "Passagens de ônibus e metrô são preços administrados com reajustes anuais. "
         "O choque no diesel não se traduz em alta imediata de tarifa."),
        ("Fricção Contratual", "Contratos de frete rodoviário têm cláusulas de reajuste com periodicidade "
         "trimestral ou semestral, absorvendo o choque gradualmente."),
        ("Diluição na Cadeia", "O diesel afeta custos industriais que se diluem em vários "
         "produtos (alimentos, manufaturas) sem aparecer diretamente no grupo de Transportes do IPCA."),
        ("Padrão Confirmado", "O diesel atua como choque de custo logístico puro e temporário — "
         "consistente com a literatura de pass-through setorial."),
    ]
    y_r2 = Inches(2.2)
    for titulo, desc in razoes:
        add_text_box(slide, titulo,
                     left=Inches(8.4), top=y_r2,
                     width=Inches(4.55), height=Inches(0.35),
                     font_size=Pt(10.5), bold=True, color=GOLD)
        add_text_box(slide, desc,
                     left=Inches(8.4), top=y_r2 + Inches(0.35),
                     width=Inches(4.55), height=Inches(0.7),
                     font_size=Pt(9.5), color=LIGHT_GRAY)
        y_r2 += Inches(1.15)

    # Box base
    add_rect(slide, left=Inches(0.3), top=Inches(6.6),
             width=Inches(12.7), height=Inches(0.85),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(0.3), top=Inches(6.6),
             width=Pt(4), height=Inches(0.85), fill_color=ACCENT)
    add_text_box(slide,
                 "Achado: O diesel afeta a inflação de transportes no curtíssimo prazo (h=0 e h=1), "
                 "mas o efeito é absorvido pela rigidez regulatória e contratual do setor de transporte brasileiro.",
                 left=Inches(0.55), top=Inches(6.65),
                 width=Inches(12.15), height=Inches(0.75),
                 font_size=Pt(10.5), italic=True, color=WHITE)


def slide_10_mediacao(prs):
    """Slide 10: Canal dos Combustíveis como Mediador"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 10)

    add_slide_title(slide, "Combustíveis como Canal Mediador",
                    "O petróleo afeta a inflação exclusivamente por dentro dos postos de gasolina")

    # Dois cenários visuais
    # Cenário 1: sem bloquear o canal
    add_rect(slide, left=Inches(0.3), top=Inches(1.6),
             width=Inches(5.9), height=Inches(2.6),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(0.3), top=Inches(1.6),
             width=Inches(5.9), height=Pt(5), fill_color=ACCENT)
    add_text_box(slide, "Modelo SEM controle de combustíveis",
                 left=Inches(0.4), top=Inches(1.65),
                 width=Inches(5.7), height=Inches(0.45),
                 font_size=Pt(12), bold=True, color=ACCENT)
    add_text_box(slide,
                 "Petróleo → IPCA Transportes\n\nCoeficiente acumulado: +0,165 ***\n(significante a 1% em h=0 a h=3)",
                 left=Inches(0.4), top=Inches(2.2),
                 width=Inches(5.7), height=Inches(1.85),
                 font_size=Pt(13), bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)

    add_text_box(slide, "→",
                 left=Inches(6.3), top=Inches(2.55),
                 width=Inches(0.7), height=Inches(0.7),
                 font_size=Pt(30), bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Cenário 2: bloqueando o canal
    add_rect(slide, left=Inches(7.1), top=Inches(1.6),
             width=Inches(5.9), height=Inches(2.6),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(7.1), top=Inches(1.6),
             width=Inches(5.9), height=Pt(5), fill_color=MID_GRAY)
    add_text_box(slide, "Modelo COM controle de combustíveis",
                 left=Inches(7.2), top=Inches(1.65),
                 width=Inches(5.7), height=Inches(0.45),
                 font_size=Pt(12), bold=True, color=MID_GRAY)
    add_text_box(slide,
                 "Petróleo → IPCA Transportes\n\nCoeficiente acumulado: ≈ 0 (ns)\n(não significante — canal bloqueado)",
                 left=Inches(7.2), top=Inches(2.2),
                 width=Inches(5.7), height=Inches(1.85),
                 font_size=Pt(13), bold=True, color=RED_WARN, align=PP_ALIGN.CENTER)

    # Diagrama de mediação
    add_rect(slide, left=Inches(0.3), top=Inches(4.45),
             width=Inches(12.7), height=Inches(2.2),
             fill_color=DARK_BLUE)
    add_text_box(slide, "Diagrama de Mediação",
                 left=Inches(0.4), top=Inches(4.5),
                 width=Inches(4.0), height=Inches(0.4),
                 font_size=Pt(12), bold=True, color=GOLD)

    # Caixas do diagrama
    med_items = [
        (Inches(0.6), "BRENT"),
        (Inches(4.9), "GASOLINA\n/ DIESEL"),
        (Inches(9.2), "IPCA\nTRANSP."),
    ]
    for mx, mlabel in med_items:
        add_rect(slide, left=mx, top=Inches(5.05),
                 width=Inches(2.5), height=Inches(1.1),
                 fill_color=NAVY)
        add_text_box(slide, mlabel,
                     left=mx, top=Inches(5.15),
                     width=Inches(2.5), height=Inches(0.9),
                     font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "→ FORTE ***", left=Inches(3.2), top=Inches(5.4),
                 width=Inches(1.6), height=Inches(0.4),
                 font_size=Pt(10), bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)
    add_text_box(slide, "→ FORTE ***", left=Inches(7.5), top=Inches(5.4),
                 width=Inches(1.6), height=Inches(0.4),
                 font_size=Pt(10), bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)
    add_text_box(slide, "↘  FRACO (canal direto, não por combustíveis)",
                 left=Inches(2.3), top=Inches(5.85),
                 width=Inches(7.3), height=Inches(0.4),
                 font_size=Pt(9.5), italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)

    # Box conclusão
    add_rect(slide, left=Inches(0.3), top=Inches(6.75),
             width=Inches(12.7), height=Inches(0.85),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(0.3), top=Inches(6.75),
             width=Pt(4), height=Inches(0.85), fill_color=GOLD)
    add_text_box(slide,
                 "Prova de Mediação: Ao controlar pelos preços de combustíveis na projeção local, o coeficiente "
                 "direto do petróleo sobre a inflação cai para zero. O petróleo só gera inflação se o canal "
                 "dos combustíveis domésticos estiver aberto.",
                 left=Inches(0.55), top=Inches(6.8),
                 width=Inches(12.15), height=Inches(0.75),
                 font_size=Pt(10.5), italic=True, color=WHITE)


def slide_11_regimes(prs):
    """Slide 11: Resultados por Regimes — Efeito PPI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 11)

    add_slide_title(slide, "Análise por Regimes: O Efeito PPI",
                    "A quebra de setembro/2016 alterou estruturalmente a transmissão inflacionária")

    # Painel de dados por regime
    add_rect(slide, left=Inches(0.3), top=Inches(1.6),
             width=Inches(7.0), height=Inches(3.5),
             fill_color=DARK_BLUE)
    add_text_box(slide, "Respostas Acumuladas por Regime (Gasolina → IPCA Transportes)",
                 left=Inches(0.4), top=Inches(1.65),
                 width=Inches(6.8), height=Inches(0.4),
                 font_size=Pt(11), bold=True, color=GOLD)

    headers_r = ["Horizonte", "Pré-2016", "IC 90%", "Pós-2016", "IC 90%", "Wald p-val."]
    data_r = [
        ["h = 0",  "0,005",  "[−0,02; 0,03]", "0,139 ***", "[0,10; 0,17]",   "< 0,01"],
        ["h = 3",  "0,052",  "[0,00; 0,17]",  "0,308 ***", "[0,15; 0,47]",   "< 0,05"],
        ["h = 6",  "0,087",  "[0,00; 0,17]",  "0,465 ***", "[0,19; 0,54]",   "< 0,05"],
        ["h = 12", "n.sig.", "[cruza zero]",   "0,450 ***", "[0,20; 0,70]",   "< 0,05"],
    ]

    col_ws = [Inches(1.0), Inches(0.95), Inches(1.4), Inches(1.1), Inches(1.4), Inches(1.0)]
    rh = Inches(0.42)

    x_tbl = Inches(0.35)
    y_tbl = Inches(2.15)

    for j, hdr in enumerate(headers_r):
        x_ = x_tbl + sum(col_ws[:j])
        add_rect(slide, left=x_, top=y_tbl,
                 width=col_ws[j], height=rh, fill_color=ACCENT)
        add_text_box(slide, hdr, left=x_ + Inches(0.03), top=y_tbl + Inches(0.06),
                     width=col_ws[j] - Inches(0.05), height=rh - Inches(0.1),
                     font_size=Pt(8.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(data_r):
        y_row = y_tbl + (ri + 1) * rh
        fill_c = DARK_BLUE if ri % 2 == 0 else RGBColor(0x16, 0x27, 0x50)
        for j, cell in enumerate(row):
            x_ = x_tbl + sum(col_ws[:j])
            add_rect(slide, left=x_, top=y_row,
                     width=col_ws[j], height=rh, fill_color=fill_c)
            if j == 0:
                tc = LIGHT_GRAY
            elif j in (1, 2) and "n.sig" not in cell:
                tc = RGBColor(0x7F, 0x8C, 0x8D)
            elif j in (3, 4) and "***" in cell:
                tc = GREEN_OK
            elif "< 0" in cell:
                tc = GREEN_OK
            else:
                tc = MID_GRAY
            add_text_box(slide, cell, left=x_ + Inches(0.03), top=y_row + Inches(0.06),
                         width=col_ws[j] - Inches(0.05), height=rh - Inches(0.1),
                         font_size=Pt(8.5), color=tc, align=PP_ALIGN.CENTER)

    # Representação visual dos dois regimes
    add_rect(slide, left=Inches(0.3), top=Inches(5.25),
             width=Inches(7.0), height=Inches(0.5),
             fill_color=DARK_BLUE)
    add_text_box(slide, "Linha AZUL (Pré-2016):  coeficiente ≈ 0  |  sem significância",
                 left=Inches(0.4), top=Inches(5.28),
                 width=Inches(6.8), height=Inches(0.4),
                 font_size=Pt(10), color=RGBColor(0x5D, 0x8A, 0xC6))

    add_rect(slide, left=Inches(0.3), top=Inches(5.8),
             width=Inches(7.0), height=Inches(0.5),
             fill_color=DARK_BLUE)
    add_text_box(slide, "Linha VERMELHA (Pós-2016):  repasse imediato e robusto",
                 left=Inches(0.4), top=Inches(5.83),
                 width=Inches(6.8), height=Inches(0.4),
                 font_size=Pt(10), bold=True, color=RED_WARN)

    # Painel direito interpretação
    add_rect(slide, left=Inches(7.6), top=Inches(1.6),
             width=Inches(5.4), height=Inches(4.7),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(7.6), top=Inches(1.6),
             width=Pt(4), height=Inches(4.7), fill_color=GOLD)
    add_text_box(slide, "O que o Teste de Wald nos diz?",
                 left=Inches(7.7), top=Inches(1.65),
                 width=Inches(5.2), height=Inches(0.4),
                 font_size=Pt(12), bold=True, color=GOLD)

    interps = [
        ("H₀ Rejeitada (p < 0,05)",
         "Os coeficientes dos dois regimes são estatisticamente distintos. "
         "A mudança de 2016 não é uma flutuação amostral."),
        ("Amortecedor Removido",
         "Antes do PPI, a Petrobras absorvia o choque em seu balanço. "
         "O consumidor estava protegido da volatilidade do Brent."),
        ("Maior Vulnerabilidade",
         "Após 2016, tensões geopolíticas (OPEP, conflitos) e desvalorizações "
         "do Real se transmitem diretamente para a bomba de gasolina."),
        ("Implicação Monetária",
         "O Banco Central passou a enfrentar choques de oferta mais frequentes "
         "e imediatos, exigindo calibragem mais precisa da política de juros."),
    ]

    y_int = Inches(2.2)
    for titulo, desc in interps:
        add_text_box(slide, titulo,
                     left=Inches(7.75), top=y_int,
                     width=Inches(5.15), height=Inches(0.38),
                     font_size=Pt(10.5), bold=True, color=ACCENT)
        add_text_box(slide, desc,
                     left=Inches(7.75), top=y_int + Inches(0.38),
                     width=Inches(5.15), height=Inches(0.65),
                     font_size=Pt(9.5), color=LIGHT_GRAY)
        y_int += Inches(1.1)

    # Box base
    add_rect(slide, left=Inches(0.3), top=Inches(6.45),
             width=Inches(12.7), height=Inches(1.0),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(0.3), top=Inches(6.45),
             width=Pt(4), height=Inches(1.0), fill_color=RED_WARN)
    add_text_box(slide,
                 "Este é o achado mais relevante do TCC: A política de preços da Petrobras "
                 "é o principal determinante da intensidade do repasse inflacionário no Brasil. "
                 "O PPI transformou a transmissão de bloqueada para imediata e robusta.",
                 left=Inches(0.55), top=Inches(6.52),
                 width=Inches(12.15), height=Inches(0.85),
                 font_size=Pt(11), bold=True, color=WHITE)


def slide_12_robustez(prs):
    """Slide 12: Robustez e Identificação Causal"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 12)

    add_slide_title(slide, "Robustez: Identificação Causal",
                    "Variáveis Instrumentais com Oil Supply News Shock (Känzig, 2021)")

    # Três etapas em cascata
    steps = [
        {
            "num": "1",
            "titulo": "O Problema: Endogeneidade Potencial",
            "desc": ("O preço do petróleo pode subir porque a economia global ou brasileira está aquecida. "
                     "Isso contamina a estimação OLS com viés de simultaneidade e variável omitida."),
            "cor": RED_WARN,
            "icon": "⚠️",
        },
        {
            "num": "2",
            "titulo": "Tentativa: LP-IV com Índice Kilian",
            "desc": ("O Índice Kilian de demanda global falhou como instrumento contemporâneo. "
                     "Estatística F do 1.º estágio < 1 (instrumento fraco por Staiger & Stock, 1997). "
                     "Solução adotada: usar Kilian como controle de demanda global no modelo OLS."),
            "cor": GOLD,
            "icon": "🔍",
        },
        {
            "num": "3",
            "titulo": "Solução: LP-IV com Oil Supply News Shock",
            "desc": ("Instrumento externo baseado em surpresas de oferta da OPEP e revisões de capacidade "
                     "(Känzig, 2021). Ortogonal à demanda brasileira.  "
                     "F do 1.º estágio ≈ 120  →  instrumento FORTE (Staiger & Stock: F > 10)."),
            "cor": GREEN_OK,
            "icon": "✅",
        },
    ]

    for i, s in enumerate(steps):
        y_box = Inches(1.6) + i * Inches(1.55)
        add_rect(slide, left=Inches(0.3), top=y_box,
                 width=Inches(12.7), height=Inches(1.4),
                 fill_color=DARK_BLUE)
        add_rect(slide, left=Inches(0.3), top=y_box,
                 width=Inches(0.65), height=Inches(1.4),
                 fill_color=s["cor"])
        add_text_box(slide, s["num"],
                     left=Inches(0.3), top=y_box + Inches(0.4),
                     width=Inches(0.65), height=Inches(0.5),
                     font_size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(slide, f"{s['icon']}  {s['titulo']}",
                     left=Inches(1.1), top=y_box + Inches(0.05),
                     width=Inches(11.5), height=Inches(0.42),
                     font_size=Pt(12), bold=True, color=s["cor"])
        add_text_box(slide, s["desc"],
                     left=Inches(1.1), top=y_box + Inches(0.5),
                     width=Inches(11.5), height=Inches(0.8),
                     font_size=Pt(10.5), color=LIGHT_GRAY)

    # Tabela de comparação LP-OLS vs LP-IV
    add_rect(slide, left=Inches(0.3), top=Inches(6.35),
             width=Inches(12.7), height=Inches(1.15),
             fill_color=DARK_BLUE)
    add_text_box(slide, "Validação: LP-OLS vs. LP-IV com OSNS  (Diesel → IPCA Transportes)",
                 left=Inches(0.4), top=Inches(6.38),
                 width=Inches(12.3), height=Inches(0.35),
                 font_size=Pt(11), bold=True, color=GOLD)

    comp_data = [
        ("h=0", "0,115 ***", "0,108 ***"),
        ("h=1", "0,166 ***", "0,161 ***"),
        ("h=2", "0,150 *",   "0,142 *"),
        ("h=6", "−0,044",    "−0,039"),
    ]

    x_comp = Inches(0.5)
    y_comp = Inches(6.78)
    for h_v, ols_v, iv_v in comp_data:
        add_text_box(slide, h_v, left=x_comp, top=y_comp,
                     width=Inches(0.6), height=Inches(0.32),
                     font_size=Pt(9), bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_text_box(slide, f"OLS: {ols_v}", left=x_comp + Inches(0.65), top=y_comp,
                     width=Inches(1.4), height=Inches(0.32),
                     font_size=Pt(9), bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)
        add_text_box(slide, f"IV: {iv_v}", left=x_comp + Inches(2.1), top=y_comp,
                     width=Inches(1.4), height=Inches(0.32),
                     font_size=Pt(9), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        x_comp += Inches(3.2)

    add_text_box(slide, "Coeficientes OLS ≈ IV  →  viés de endogeneidade desprezível  →  modelo principal VALIDADO",
                 left=Inches(0.4), top=Inches(7.12),
                 width=Inches(12.5), height=Inches(0.3),
                 font_size=Pt(10), bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)


def slide_13_conclusoes(prs):
    """Slide 13: Conclusões da Pesquisa"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 13)

    add_slide_title(slide, "Conclusões da Pesquisa",
                    "Três achados centrais, robustos e coerentes com a teoria econômica")

    conclusoes = [
        {
            "num": "01",
            "titulo": "Canal de Custo Ativo e Setorialmente Concentrado",
            "desc": ("O petróleo transmite-se de forma robusta aos combustíveis e destes ao IPCA Transportes "
                     "(pass-through de ~43% em 12 meses, t > 3,6 em todos os horizontes). O IPCA Geral "
                     "apresenta impacto modesto e transitório, típico de choque de custo setorial."),
            "cor": ACCENT,
            "stat": "Pass-through\n~43%",
        },
        {
            "num": "02",
            "titulo": "Assimetria Temporal entre Gasolina e Diesel",
            "desc": ("A gasolina gera repasse cumulativo e persistente no médio prazo (h=12 ainda significante). "
                     "O diesel produz choque logístico concentrado no curtíssimo prazo (h=0 e h=1), dissipando-se "
                     "por causa da rigidez regulatória das tarifas de transporte público."),
            "cor": GOLD,
            "stat": "Diesel: 2 meses\nGasolina: 12+ meses",
        },
        {
            "num": "03",
            "titulo": "A Política da Petrobras Determina a Velocidade do Repasse",
            "desc": ("O PPI de 2016 foi uma mudança estrutural comprovada estatisticamente (Wald p < 0,05). "
                     "No regime anterior, o repasse estava bloqueado. No regime atual, é imediato e robusto, "
                     "tornando a inflação brasileira mais sensível às tensões geopolíticas globais."),
            "cor": GREEN_OK,
            "stat": "Wald\np < 0,05",
        },
    ]

    for i, c in enumerate(conclusoes):
        y_box = Inches(1.6) + i * Inches(1.7)
        add_rect(slide, left=Inches(0.3), top=y_box,
                 width=Inches(11.2), height=Inches(1.55),
                 fill_color=DARK_BLUE)
        add_rect(slide, left=Inches(0.3), top=y_box,
                 width=Inches(0.65), height=Inches(1.55),
                 fill_color=c["cor"])
        add_text_box(slide, c["num"],
                     left=Inches(0.3), top=y_box + Inches(0.45),
                     width=Inches(0.65), height=Inches(0.55),
                     font_size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(slide, c["titulo"],
                     left=Inches(1.1), top=y_box + Inches(0.1),
                     width=Inches(9.0), height=Inches(0.4),
                     font_size=Pt(12), bold=True, color=c["cor"])
        add_text_box(slide, c["desc"],
                     left=Inches(1.1), top=y_box + Inches(0.52),
                     width=Inches(9.0), height=Inches(0.95),
                     font_size=Pt(10.5), color=LIGHT_GRAY)

        # Stat box
        add_rect(slide, left=Inches(11.65), top=y_box,
                 width=Inches(1.35), height=Inches(1.55),
                 fill_color=c["cor"])
        add_text_box(slide, c["stat"],
                     left=Inches(11.65), top=y_box + Inches(0.25),
                     width=Inches(1.35), height=Inches(1.0),
                     font_size=Pt(12), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Box avaliação
    add_rect(slide, left=Inches(0.3), top=Inches(6.75),
             width=Inches(12.7), height=Inches(0.7),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(0.3), top=Inches(6.75),
             width=Pt(4), height=Inches(0.7), fill_color=GOLD)
    add_text_box(slide,
                 "Todos os achados são consistentes entre si e com a literatura macroeconômica. "
                 "O trabalho contribui com evidência empírica inédita sobre o papel moderador "
                 "da estatal de energia na dinâmica inflacionária do Brasil (2003–2026).",
                 left=Inches(0.55), top=Inches(6.8),
                 width=Inches(12.15), height=Inches(0.6),
                 font_size=Pt(10), italic=True, color=LIGHT_GRAY)


def slide_14_politica(prs):
    """Slide 14: Implicações para Política Monetária e Encerramento"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    add_accent_bar(slide)
    add_footer(slide, 14)

    add_slide_title(slide, "Implicações para a Política Econômica",
                    "O que os resultados significam para o Banco Central e para os formuladores de política?")

    impl = [
        {
            "icon": "🏦",
            "titulo": "Desafio para o Banco Central",
            "desc": ("Sob o regime PPI, choques de oferta globais geram pressão inflacionária imediata sobre "
                     "o IPCA Transportes. O BCB enfrenta o dilema clássico: reagir com juros arriscando contrair "
                     "a economia, ou tolerar temporariamente a alta sabendo que o efeito sobre o IPCA Geral "
                     "é transitório (dissipação em 5–6 meses)."),
            "cor": ACCENT,
        },
        {
            "icon": "⚡",
            "titulo": "Risco de Sobre-Reação Monetária",
            "desc": ("Como o IPCA Geral perde significância após o 6.º mês, uma alta de Selic defensiva "
                     "para conter a inflação de transporte pode gerar custo de atividade desnecessário. "
                     "Os resultados sugerem cautela na resposta monetária a choques de custo energético "
                     "de natureza temporária."),
            "cor": GOLD,
        },
        {
            "icon": "🏭",
            "titulo": "O Papel Regulatório da Petrobras",
            "desc": ("Os formuladores de política precisam reconhecer que a estrutura de precificação da "
                     "Petrobras é um instrumento implícito de política de estabilização inflacionária. "
                     "A escolha entre mercado livre (PPI) e preços administrados tem impacto direto e "
                     "mensurável sobre a inflação setorial de transportes."),
            "cor": GREEN_OK,
        },
    ]

    for i, it in enumerate(impl):
        y_box = Inches(1.6) + i * Inches(1.55)
        add_rect(slide, left=Inches(0.3), top=y_box,
                 width=Inches(12.7), height=Inches(1.4),
                 fill_color=DARK_BLUE)
        add_rect(slide, left=Inches(0.3), top=y_box,
                 width=Inches(0.65), height=Inches(1.4),
                 fill_color=it["cor"])
        add_text_box(slide, it["icon"],
                     left=Inches(0.3), top=y_box + Inches(0.35),
                     width=Inches(0.65), height=Inches(0.55),
                     font_size=Pt(20), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(slide, it["titulo"],
                     left=Inches(1.1), top=y_box + Inches(0.08),
                     width=Inches(11.5), height=Inches(0.4),
                     font_size=Pt(12), bold=True, color=it["cor"])
        add_text_box(slide, it["desc"],
                     left=Inches(1.1), top=y_box + Inches(0.5),
                     width=Inches(11.5), height=Inches(0.8),
                     font_size=Pt(10.5), color=LIGHT_GRAY)

    # Agradecimento
    add_rect(slide, left=Inches(0.3), top=Inches(6.4),
             width=Inches(12.7), height=Inches(1.2),
             fill_color=DARK_BLUE)
    add_rect(slide, left=Inches(0.3), top=Inches(6.4),
             width=Inches(12.7), height=Pt(3), fill_color=GOLD)
    add_text_box(slide, "Obrigado!",
                 left=Inches(0.4), top=Inches(6.45),
                 width=Inches(12.3), height=Inches(0.55),
                 font_size=Pt(22), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide,
                 "Pedro Franck Minella   |   pedrofrankminella@gmail.com   |   Ibmec-DF, 2026   "
                 "|   Orientador: Prof. Silvio Costa",
                 left=Inches(0.4), top=Inches(7.0),
                 width=Inches(12.3), height=Inches(0.4),
                 font_size=Pt(10), color=MID_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Gerando Slide 1: Capa...")
    slide_01_capa(prs)
    print("Gerando Slide 2: Mecanismo de Transmissão...")
    slide_02_mecanismo(prs)
    print("Gerando Slide 3: Petrobras Institucional...")
    slide_03_petrobras(prs)
    print("Gerando Slide 4: Perguntas de Pesquisa...")
    slide_04_perguntas(prs)
    print("Gerando Slide 5: Metodologia LP...")
    slide_05_metodologia(prs)
    print("Gerando Slide 6: Estratégia Empírica...")
    slide_06_estrategia(prs)
    print("Gerando Slide 7: Petróleo → Combustíveis...")
    slide_07_combustiveis(prs)
    print("Gerando Slide 8: Gasolina → IPCA Transportes vs. Geral...")
    slide_08_gasolina_ipca(prs)
    print("Gerando Slide 9: Diesel — Impacto Temporário...")
    slide_09_diesel(prs)
    print("Gerando Slide 10: Canal de Mediação...")
    slide_10_mediacao(prs)
    print("Gerando Slide 11: Regimes — Efeito PPI...")
    slide_11_regimes(prs)
    print("Gerando Slide 12: Robustez e IV...")
    slide_12_robustez(prs)
    print("Gerando Slide 13: Conclusões...")
    slide_13_conclusoes(prs)
    print("Gerando Slide 14: Política Monetária e Encerramento...")
    slide_14_politica(prs)

    output_path = "Apresentacao_TCC_Pedro_Minella.pptx"
    prs.save(output_path)
    print(f"\n✅ Apresentação salva com sucesso em: {output_path}")
    print(f"   Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
