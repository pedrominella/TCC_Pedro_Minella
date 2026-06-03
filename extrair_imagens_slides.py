# -*- coding: utf-8 -*-
"""Extrai as imagens de cada slide do NotebookLM para inspecionar"""
from pptx import Presentation
from pathlib import Path
import shutil

prs = Presentation(r'C:\Users\pedro\OneDrive\Documentos\TCC_python\Oil_Shocks_and_Brazilian_Inflation.pptx')

out_dir = Path(r'C:\Users\pedro\OneDrive\Documentos\TCC_python\inspecao_slides_notebooklm')
out_dir.mkdir(exist_ok=True)

for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            img_blob = shape.image.blob
            ext = shape.image.ext
            fname = out_dir / f'slide_{i+1:02d}.{ext}'
            with open(fname, 'wb') as f:
                f.write(img_blob)
            print(f'Slide {i+1}: salvo em {fname}')

print('\nConcluido!')
